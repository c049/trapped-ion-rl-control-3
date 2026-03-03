#!/usr/bin/env python3
"""Build a detailed Sydney Uni weekly report PPT.

Design goals from meeting notes:
- Motivation -> Setup/Method -> Results -> Outlook story.
- Picture + words first, no code snippets.
- Explain RL loop clearly (episode/epoch/policy update/PPO).
- Emphasize characteristic-function reward and sampling-point evolution.
- Include all three state GIFs (Cat/GKP/Binomial).
- Summarize robust training (quasi-static + stochastic) with latest sweep results.
"""

import re
import shutil
from datetime import datetime
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, List, Optional, Tuple

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PPT_DIR = ROOT / "ppt"
ASSET_DIR = PPT_DIR / "assets"

CAT_OUT = ROOT / "examples" / "trapped_ion_cat" / "outputs"
GKP_OUT = ROOT / "examples" / "trapped_ion_gkp" / "outputs"
BIN_OUT = ROOT / "examples" / "trapped_ion_binomial" / "outputs"
BIN_CHECKPOINT = ROOT / "examples" / "trapped_ion_binomial" / "checkpoint"
BIN_PRE_ROBUST_OUT = BIN_OUT / "archive_20260225_144627_pre_new_run"
BIN_SWEEP_ROOT = ROOT / "examples" / "trapped_ion_binomial" / "penalty_sweep"

STATIC_SWEEP = BIN_SWEEP_ROOT / "penalty_sweep_20260225_195742"
STOCH_SWEEP_COARSE = BIN_SWEEP_ROOT / "stochastic_penalty_sweep_20260226_023916"
STOCH_SWEEP_REFINE = BIN_SWEEP_ROOT / "stoch_validate_uniform_weight_local_20260226_170726"

PPT_PATH_DETAILED = PPT_DIR / "Sydney_Uni_Weekly_Report_2026-03-03_detailed.pptx"
PPT_PATH_ALIAS = PPT_DIR / "Sydney_Uni_Weekly_Report_2026-03-03.pptx"


@dataclass
class SweepEntry:
    mode: str
    tag: str
    penalty: float
    score: float
    f_nom: float
    f_rob: float
    penalty_term: float
    run_dir: Path


@dataclass
class DephasingStats:
    mean_robust: float
    mean_baseline: float
    mean_gain: float
    min_gain: float
    max_gain: float
    zero_gain: float


@dataclass
class HyperParamSummary:
    cat_num_epochs: str
    cat_policy_updates: str
    cat_n_steps: str
    cat_n_segments: str
    cat_stage_points: str
    cat_stage_epochs: str
    gkp_num_epochs: str
    gkp_policy_updates: str
    gkp_n_steps: str
    gkp_n_segments: str
    gkp_stage_points: str
    gkp_stage_epochs: str
    bin_num_epochs: str
    bin_policy_updates: str
    bin_n_steps: str
    bin_n_segments: str
    bin_stage_points: str
    bin_stage_epochs: str


@dataclass
class ConstantAmpRun:
    tag: str
    fidelity: float
    robust_enabled: Optional[bool]
    amp_refine_enabled: Optional[bool]
    amp_full_enabled: Optional[bool]
    path: Path


@dataclass
class Metrics:
    cat_fid: float
    gkp_fid: float
    bin_pre_fid: float
    bin_pre_score: float
    bin_pre_f_rob: float
    bin_pre_run_dir: Path
    static_entries: List[SweepEntry]
    stoch_entries: List[SweepEntry]
    best_static: SweepEntry
    best_stoch: SweepEntry
    static_stats: DephasingStats
    stoch_stats: DephasingStats
    hyper: HyperParamSummary
    nonrobust_history: List[Tuple[str, float]]
    best_nonrobust_tag: str
    best_nonrobust_fid: float
    checkpoint_nonrobust_fid: float
    baseline_fid_at_zero_detuning: float
    checkpoint_amp_r_std: float
    checkpoint_amp_b_std: float
    constamp_nonrobust: Optional[ConstantAmpRun]
    constamp_robust: Optional[ConstantAmpRun]


def _font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"
        if bold
        else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/dejavu/DejaVuSans-Bold.ttf"
        if bold
        else "/usr/share/fonts/dejavu/DejaVuSans.ttf",
    ]
    for p in candidates:
        fp = Path(p)
        if fp.exists():
            return ImageFont.truetype(str(fp), size=size)
    return ImageFont.load_default()


def _rr(
    draw: ImageDraw.ImageDraw,
    xy: Tuple[int, int, int, int],
    radius: int,
    fill,
    outline,
    width: int = 3,
) -> None:
    draw.rounded_rectangle(xy, radius=radius, fill=fill, outline=outline, width=width)


def _arrow(
    draw: ImageDraw.ImageDraw,
    p1: Tuple[int, int],
    p2: Tuple[int, int],
    color,
    width: int = 5,
    head: int = 18,
) -> None:
    draw.line([p1, p2], fill=color, width=width)
    x1, y1 = p1
    x2, y2 = p2
    dx = x2 - x1
    dy = y2 - y1
    norm = (dx * dx + dy * dy) ** 0.5 or 1.0
    ux = dx / norm
    uy = dy / norm
    px = -uy
    py = ux
    hx = x2 - int(head * ux)
    hy = y2 - int(head * uy)
    poly = [
        (x2, y2),
        (hx + int(0.6 * head * px), hy + int(0.6 * head * py)),
        (hx - int(0.6 * head * px), hy - int(0.6 * head * py)),
    ]
    draw.polygon(poly, fill=color)


def _wrap_text_lines(
    draw: ImageDraw.ImageDraw,
    text: str,
    font: ImageFont.FreeTypeFont,
    max_width: int,
) -> List[str]:
    """Wrap text to fit pixel width for PIL drawing while preserving paragraphs."""
    out: List[str] = []
    for para in text.split("\n"):
        words = para.split()
        if not words:
            out.append("")
            continue
        current = words[0]
        for w in words[1:]:
            candidate = current + " " + w
            bb = draw.textbbox((0, 0), candidate, font=font)
            width = bb[2] - bb[0]
            if width <= max_width:
                current = candidate
            else:
                out.append(current)
                current = w
        out.append(current)
    return out


def _wrap_text(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.FreeTypeFont, max_width: int) -> str:
    return "\n".join(_wrap_text_lines(draw, text, font, max_width))


def _fit_wrapped_text(
    draw: ImageDraw.ImageDraw,
    text: str,
    max_width: int,
    max_height: int,
    start_size: int,
    min_size: int = 14,
    bold: bool = False,
    spacing: int = 6,
) -> Tuple[ImageFont.FreeTypeFont, str]:
    """Find a wrapped-text font size that fits within a box."""
    if max_width <= 2 or max_height <= 2:
        font = _font(min_size, bold=bold)
        return font, text

    for size in range(start_size, min_size - 1, -1):
        font = _font(size, bold=bold)
        wrapped = "\n".join(_wrap_text_lines(draw, text, font, max_width))
        bb = draw.multiline_textbbox((0, 0), wrapped, font=font, spacing=spacing)
        tw = bb[2] - bb[0]
        th = bb[3] - bb[1]
        if tw <= max_width and th <= max_height:
            return font, wrapped

    font = _font(min_size, bold=bold)
    wrapped = "\n".join(_wrap_text_lines(draw, text, font, max_width))
    return font, wrapped


def _draw_wrapped_text(
    draw: ImageDraw.ImageDraw,
    xy: Tuple[int, int],
    text: str,
    font: ImageFont.FreeTypeFont,
    max_width: int,
    fill,
    spacing: int = 6,
) -> None:
    wrapped = _wrap_text(draw, text, font, max_width)
    draw.multiline_text(xy, wrapped, font=font, fill=fill, spacing=spacing)


def _draw_text_in_box(
    draw: ImageDraw.ImageDraw,
    box: Tuple[int, int, int, int],
    text: str,
    *,
    start_size: int,
    min_size: int = 14,
    bold: bool = False,
    fill=(30, 46, 72),
    spacing: int = 6,
    padding: int = 0,
    align: str = "left",
    valign: str = "top",
) -> Tuple[int, int]:
    """Draw wrapped text inside a box with automatic font down-scaling."""
    x1, y1, x2, y2 = box
    inner_x1 = x1 + padding
    inner_y1 = y1 + padding
    inner_x2 = x2 - padding
    inner_y2 = y2 - padding
    max_width = max(2, inner_x2 - inner_x1)
    max_height = max(2, inner_y2 - inner_y1)
    font, wrapped = _fit_wrapped_text(
        draw,
        text,
        max_width=max_width,
        max_height=max_height,
        start_size=start_size,
        min_size=min_size,
        bold=bold,
        spacing=spacing,
    )
    bb = draw.multiline_textbbox((0, 0), wrapped, font=font, spacing=spacing, align=align)
    tw = bb[2] - bb[0]
    th = bb[3] - bb[1]

    if align == "center":
        tx = inner_x1 + max(0, (max_width - tw) // 2)
    elif align == "right":
        tx = inner_x1 + max(0, max_width - tw)
    else:
        tx = inner_x1

    if valign == "middle":
        ty = inner_y1 + max(0, (max_height - th) // 2)
    elif valign == "bottom":
        ty = inner_y1 + max(0, max_height - th)
    else:
        ty = inner_y1

    draw.multiline_text((tx, ty), wrapped, font=font, fill=fill, spacing=spacing, align=align)
    return tw, th


def _safe_float_text(path: Path, default: float = 0.0) -> float:
    try:
        return float(path.read_text(encoding="utf-8").strip())
    except Exception:
        return default


def _read_score_file(path: Path) -> Dict[str, float]:
    out = {}
    txt = path.read_text(encoding="utf-8", errors="ignore").strip()
    for item in txt.split(","):
        if "=" in item:
            k, v = item.split("=", 1)
            try:
                out[k.strip()] = float(v.strip())
            except Exception:
                continue
    return out


def _parse_penalty_token(token: str) -> float:
    t = token.strip()
    if t == "":
        return 0.0
    if t == "0":
        return 0.0
    if "p" in t:
        t = t.replace("p", ".")
    return float(t)


def _collect_sweep_entries(root: Path, prefix: str, mode: str) -> List[SweepEntry]:
    entries: List[SweepEntry] = []
    if not root.exists():
        return entries
    for run_dir in sorted(root.iterdir()):
        if not run_dir.is_dir() or not run_dir.name.startswith(prefix):
            continue
        score_file = run_dir / "final_robust_score.txt"
        if not score_file.exists():
            continue
        vals = _read_score_file(score_file)
        token = run_dir.name[len(prefix) :]
        try:
            penalty = _parse_penalty_token(token)
        except Exception:
            continue
        entries.append(
            SweepEntry(
                mode=mode,
                tag=run_dir.name,
                penalty=penalty,
                score=vals.get("score", 0.0),
                f_nom=vals.get("f_nom", 0.0),
                f_rob=vals.get("f_rob", 0.0),
                penalty_term=vals.get("penalty", 0.0),
                run_dir=run_dir,
            )
        )
    return entries


def _best_by_score(entries: List[SweepEntry]) -> SweepEntry:
    if not entries:
        raise RuntimeError("No sweep entries found.")
    return sorted(entries, key=lambda x: x.score, reverse=True)[0]


def _best_per_penalty(entries: List[SweepEntry]) -> List[SweepEntry]:
    keep: Dict[float, SweepEntry] = {}
    for e in entries:
        if e.penalty not in keep or e.score > keep[e.penalty].score:
            keep[e.penalty] = e
    return [keep[k] for k in sorted(keep)]


def _read_dephasing_stats(run_dir: Path) -> DephasingStats:
    csv_path = run_dir / "dephasing_compare.csv"
    if not csv_path.exists():
        return DephasingStats(0.0, 0.0, 0.0, 0.0, 0.0, 0.0)
    df = pd.read_csv(csv_path)
    if not {"robust_fidelity", "baseline_fidelity", "detuning_frac"}.issubset(df.columns):
        return DephasingStats(0.0, 0.0, 0.0, 0.0, 0.0, 0.0)
    gain = df["robust_fidelity"] - df["baseline_fidelity"]
    zero_idx = int((df["detuning_frac"].abs()).idxmin())
    return DephasingStats(
        mean_robust=float(df["robust_fidelity"].mean()),
        mean_baseline=float(df["baseline_fidelity"].mean()),
        mean_gain=float(gain.mean()),
        min_gain=float(gain.min()),
        max_gain=float(gain.max()),
        zero_gain=float(gain.iloc[zero_idx]),
    )


def _read_nonrobust_history(log_dir: Path) -> List[Tuple[str, float]]:
    """Parse non-robust (no-noise/dephase-off) final fidelities from client logs."""
    out: List[Tuple[str, float]] = []
    if not log_dir.exists():
        return out
    for p in sorted(log_dir.glob("client_*.log")):
        txt = p.read_text(encoding="utf-8", errors="ignore")
        m_fid = re.search(r"Final fidelity\s+([0-9]*\.?[0-9]+)", txt)
        if not m_fid:
            continue
        fid = float(m_fid.group(1))
        m_rob = re.search(r"Robust dephasing:\s+enabled=(True|False)", txt)
        if m_rob:
            if m_rob.group(1) == "False":
                out.append((p.stem.replace("client_", ""), fid))
            continue
        # Early logs before robust flags existed: treat as non-robust runs.
        out.append((p.stem.replace("client_", ""), fid))
    return out


def _read_checkpoint_nonrobust_best(checkpoint_dir: Path) -> Optional[Tuple[str, float]]:
    p = checkpoint_dir / "final_fidelity_best.txt"
    if not p.exists():
        return None
    try:
        fid = float(p.read_text(encoding="utf-8", errors="ignore").strip())
    except Exception:
        return None
    ts = datetime.fromtimestamp(p.stat().st_mtime).strftime("%Y%m%d_%H%M%S")
    return (f"checkpoint_{ts}", fid)


def _read_baseline_fid_at_zero(run_dir: Path) -> float:
    csv_path = run_dir / "dephasing_compare.csv"
    if not csv_path.exists():
        return 0.0
    df = pd.read_csv(csv_path)
    if not {"detuning_frac", "baseline_fidelity"}.issubset(df.columns):
        return 0.0
    idx = int((df["detuning_frac"].abs()).idxmin())
    return float(df["baseline_fidelity"].iloc[idx])


def _pulse_amp_std(npz_path: Path) -> Tuple[float, float]:
    if not npz_path.exists():
        return 0.0, 0.0
    try:
        data = np.load(npz_path)
        amp_r = np.asarray(data["amp_r"], dtype=float)
        amp_b = np.asarray(data["amp_b"], dtype=float)
        return float(np.std(amp_r)), float(np.std(amp_b))
    except Exception:
        return 0.0, 0.0


def _iter_candidate_logs() -> List[Path]:
    roots = [
        BIN_OUT / "logs",
        BIN_SWEEP_ROOT,
    ]
    out: List[Path] = []
    for root in roots:
        if not root.exists():
            continue
        out.extend(sorted(root.rglob("*.log")))
    return out


def _parse_constant_amp_run(log_path: Path) -> Optional[ConstantAmpRun]:
    txt = log_path.read_text(encoding="utf-8", errors="ignore")
    m_fid = re.search(r"Final fidelity\s+([0-9]*\.?[0-9]+)", txt)
    if not m_fid:
        return None
    fid = float(m_fid.group(1))

    m_rob = re.search(r"Robust dephasing:\s+enabled=(True|False)", txt)
    robust_enabled: Optional[bool]
    if m_rob:
        robust_enabled = m_rob.group(1) == "True"
    else:
        robust_enabled = None

    m_refine = re.search(r"Final refinement setup:.*?amp_opt=(True|False)", txt)
    amp_refine_enabled: Optional[bool]
    if m_refine:
        amp_refine_enabled = m_refine.group(1) == "True"
    else:
        amp_refine_enabled = None

    m_full = re.search(r"Full-step refinement:.*?amp_opt=(True|False)", txt)
    amp_full_enabled: Optional[bool]
    if m_full:
        amp_full_enabled = m_full.group(1) == "True"
    else:
        amp_full_enabled = None

    # Strict constant-amplitude criterion:
    # 1) final refinement does not optimize amplitudes (when field exists)
    # 2) full-step refinement does not optimize amplitudes (when field exists)
    # 3) runtime amplitude stats show std=0 for amp_r and amp_b (if present)
    if amp_refine_enabled is True or amp_full_enabled is True:
        return None

    amp_ok = True
    m_amp_r = re.findall(
        r"amp_r stats: mean=([0-9\.\-]+)\s+std=([0-9\.\-]+)\s+min=([0-9\.\-]+)\s+max=([0-9\.\-]+)",
        txt,
    )
    m_amp_b = re.findall(
        r"amp_b stats: mean=([0-9\.\-]+)\s+std=([0-9\.\-]+)\s+min=([0-9\.\-]+)\s+max=([0-9\.\-]+)",
        txt,
    )
    if m_amp_r and m_amp_b:
        try:
            amp_r_std = float(m_amp_r[-1][1])
            amp_b_std = float(m_amp_b[-1][1])
            amp_ok = (abs(amp_r_std) < 1.0e-10) and (abs(amp_b_std) < 1.0e-10)
        except Exception:
            amp_ok = False
    if not amp_ok:
        return None

    tag = log_path.stem.replace("client_", "")
    return ConstantAmpRun(
        tag=tag,
        fidelity=fid,
        robust_enabled=robust_enabled,
        amp_refine_enabled=amp_refine_enabled,
        amp_full_enabled=amp_full_enabled,
        path=log_path,
    )


def _read_constant_amp_evidence() -> Tuple[Optional[ConstantAmpRun], Optional[ConstantAmpRun]]:
    candidates: List[ConstantAmpRun] = []
    for p in _iter_candidate_logs():
        run = _parse_constant_amp_run(p)
        if run is not None:
            candidates.append(run)
    if not candidates:
        return None, None

    nonrob = [r for r in candidates if r.robust_enabled is False or r.robust_enabled is None]
    rob = [r for r in candidates if r.robust_enabled is True]

    best_nonrob = max(nonrob, key=lambda r: r.fidelity) if nonrob else None
    best_rob = max(rob, key=lambda r: r.fidelity) if rob else None
    return best_nonrob, best_rob


def _extract_env_default(path: Path, key: str, fallback: str = "-") -> str:
    if not path.exists():
        return fallback
    txt = path.read_text(encoding="utf-8", errors="ignore")
    patterns = [
        r'os\.environ\.get\("%s",\s*"([^"]+)"\)' % re.escape(key),
        r"os\.environ\.get\('%s',\s*'([^']+)'\)" % re.escape(key),
    ]
    for pat in patterns:
        m = re.search(pat, txt)
        if m:
            return m.group(1)
    return fallback


def _extract_literal_assignment(path: Path, key: str, fallback: str = "-") -> str:
    if not path.exists():
        return fallback
    txt = path.read_text(encoding="utf-8", errors="ignore")
    m = re.search(r"^\s*%s\s*=\s*([0-9\.]+)\s*$" % re.escape(key), txt, re.MULTILINE)
    if m:
        return m.group(1)
    return fallback


def _read_hyperparams() -> HyperParamSummary:
    cat_server = ROOT / "examples" / "trapped_ion_cat" / "trapped_ion_cat_training_server.py"
    cat_client = ROOT / "examples" / "trapped_ion_cat" / "trapped_ion_cat_client.py"
    gkp_server = ROOT / "examples" / "trapped_ion_gkp" / "trapped_ion_gkp_training_server.py"
    gkp_client = ROOT / "examples" / "trapped_ion_gkp" / "trapped_ion_gkp_client.py"
    bin_server = ROOT / "examples" / "trapped_ion_binomial" / "trapped_ion_binomial_training_server.py"
    bin_client = ROOT / "examples" / "trapped_ion_binomial" / "trapped_ion_binomial_client.py"

    cat_s1_epoch = _extract_env_default(cat_client, "TRAIN_STAGE1_EPOCHS", "120")
    cat_s2_epoch = _extract_env_default(cat_client, "TRAIN_STAGE2_EPOCHS", "240")
    gkp_s1_epoch = _extract_env_default(gkp_client, "TRAIN_STAGE1_EPOCHS", "0")
    gkp_s2_epoch = _extract_env_default(gkp_client, "TRAIN_STAGE2_EPOCHS", "180")
    bin_s1_epoch = _extract_env_default(bin_client, "TRAIN_STAGE1_EPOCHS", "120")
    bin_s2_epoch = _extract_env_default(bin_client, "TRAIN_STAGE2_EPOCHS", "240")

    gkp_stage_epoch_text = (
        "S1=off, S2<%s" % gkp_s2_epoch if gkp_s1_epoch == "0" else "S1<%s, S2<%s" % (gkp_s1_epoch, gkp_s2_epoch)
    )

    return HyperParamSummary(
        cat_num_epochs=_extract_env_default(cat_server, "NUM_EPOCHS", "300"),
        cat_policy_updates=_extract_env_default(cat_server, "PPO_NUM_POLICY_UPDATES", "20"),
        cat_n_steps=_extract_literal_assignment(cat_client, "N_STEPS", "120"),
        cat_n_segments=_extract_literal_assignment(cat_client, "N_SEGMENTS", "60"),
        cat_stage_points=(
            "S1=%s, S2=%s, S3=%s"
            % (
                _extract_env_default(cat_client, "TRAIN_POINTS_STAGE1", "120"),
                _extract_env_default(cat_client, "TRAIN_POINTS_STAGE2", "240"),
                _extract_env_default(cat_client, "TRAIN_POINTS_STAGE3", "960"),
            )
        ),
        cat_stage_epochs=(
            "S1<%s, S2<%s"
            % (
                cat_s1_epoch,
                cat_s2_epoch,
            )
        ),
        gkp_num_epochs=_extract_env_default(gkp_server, "NUM_EPOCHS", "2000"),
        gkp_policy_updates=_extract_env_default(gkp_server, "PPO_NUM_POLICY_UPDATES", "10"),
        gkp_n_steps=_extract_env_default(gkp_server, "N_STEPS", "120"),
        gkp_n_segments=_extract_env_default(gkp_server, "N_SEGMENTS", "60"),
        gkp_stage_points=(
            "S1=%s, S2=%s, S3=%s"
            % (
                _extract_env_default(gkp_client, "TRAIN_POINTS_STAGE1", "120"),
                _extract_env_default(gkp_client, "TRAIN_POINTS_STAGE2", "240"),
                _extract_env_default(gkp_client, "TRAIN_POINTS_STAGE3", "960"),
            )
        ),
        gkp_stage_epochs=(
            gkp_stage_epoch_text
        ),
        bin_num_epochs=_extract_env_default(bin_server, "NUM_EPOCHS", "300"),
        bin_policy_updates=_extract_env_default(bin_server, "PPO_NUM_POLICY_UPDATES", "20"),
        bin_n_steps=_extract_env_default(bin_server, "N_STEPS", "120"),
        bin_n_segments=_extract_env_default(bin_server, "N_SEGMENTS", "60"),
        bin_stage_points=(
            "S1=%s, S2=%s, S3=%s"
            % (
                _extract_env_default(bin_client, "TRAIN_POINTS_STAGE1", "120"),
                _extract_env_default(bin_client, "TRAIN_POINTS_STAGE2", "240"),
                _extract_env_default(bin_client, "TRAIN_POINTS_STAGE3", "960"),
            )
        ),
        bin_stage_epochs=(
            "S1<%s, S2<%s"
            % (
                bin_s1_epoch,
                bin_s2_epoch,
            )
        ),
    )


def read_metrics() -> Metrics:
    static_entries_raw = _collect_sweep_entries(STATIC_SWEEP, "p_", "quasi-static")
    stoch_entries_raw = _collect_sweep_entries(STOCH_SWEEP_COARSE, "stoch_p_", "stochastic")
    stoch_entries_raw += _collect_sweep_entries(STOCH_SWEEP_REFINE, "stoch_p_", "stochastic")

    static_entries = _best_per_penalty(static_entries_raw)
    stoch_entries = _best_per_penalty(stoch_entries_raw)

    best_static = _best_by_score(static_entries_raw)
    best_stoch = _best_by_score(stoch_entries_raw)
    pre_vals = _read_score_file(BIN_PRE_ROBUST_OUT / "final_robust_score.txt") if (BIN_PRE_ROBUST_OUT / "final_robust_score.txt").exists() else {}
    nonrobust_hist = _read_nonrobust_history(BIN_OUT / "logs")
    checkpoint_nonrobust = _read_checkpoint_nonrobust_best(BIN_CHECKPOINT)
    if checkpoint_nonrobust is not None:
        nonrobust_hist = nonrobust_hist + [checkpoint_nonrobust]
    if nonrobust_hist:
        best_nonrobust = max(nonrobust_hist, key=lambda kv: kv[1])
    else:
        best_nonrobust = ("n/a", 0.0)

    checkpoint_amp_r_std, checkpoint_amp_b_std = _pulse_amp_std(BIN_CHECKPOINT / "final_pulses_best.npz")
    constamp_nonrobust, constamp_robust = _read_constant_amp_evidence()

    return Metrics(
        cat_fid=_safe_float_text(CAT_OUT / "final_fidelity.txt"),
        gkp_fid=_safe_float_text(GKP_OUT / "final_fidelity.txt"),
        bin_pre_fid=_safe_float_text(BIN_PRE_ROBUST_OUT / "final_fidelity.txt", default=_safe_float_text(BIN_OUT / "final_fidelity.txt")),
        bin_pre_score=float(pre_vals.get("score", 0.0)),
        bin_pre_f_rob=float(pre_vals.get("f_rob", 0.0)),
        bin_pre_run_dir=BIN_PRE_ROBUST_OUT if BIN_PRE_ROBUST_OUT.exists() else BIN_OUT,
        static_entries=static_entries,
        stoch_entries=stoch_entries,
        best_static=best_static,
        best_stoch=best_stoch,
        static_stats=_read_dephasing_stats(best_static.run_dir),
        stoch_stats=_read_dephasing_stats(best_stoch.run_dir),
        hyper=_read_hyperparams(),
        nonrobust_history=nonrobust_hist,
        best_nonrobust_tag=best_nonrobust[0],
        best_nonrobust_fid=float(best_nonrobust[1]),
        checkpoint_nonrobust_fid=float(checkpoint_nonrobust[1]) if checkpoint_nonrobust is not None else 0.0,
        baseline_fid_at_zero_detuning=_read_baseline_fid_at_zero(best_static.run_dir),
        checkpoint_amp_r_std=checkpoint_amp_r_std,
        checkpoint_amp_b_std=checkpoint_amp_b_std,
        constamp_nonrobust=constamp_nonrobust,
        constamp_robust=constamp_robust,
    )


def _make_missing_image(path: Path) -> None:
    img = Image.new("RGB", (1200, 700), (246, 248, 252))
    d = ImageDraw.Draw(img)
    d.text((80, 300), "Missing image", font=_font(64, bold=True), fill=(120, 130, 150))
    img.save(path)


def _safe_image(path: Path, fallback: Path) -> Path:
    if path.exists():
        return path
    return fallback


def build_hook_image(path: Path) -> None:
    w, h = 1600, 900
    img = Image.new("RGB", (w, h), (246, 250, 255))
    d = ImageDraw.Draw(img)
    f_title = _font(44, bold=True)
    f_sub = _font(30, bold=True)
    f_txt = _font(23)

    _draw_text_in_box(
        d,
        (60, 28, 1520, 120),
        "Main Question: Can model-free RL learn robust bosonic control pulses?",
        start_size=44,
        min_size=34,
        bold=True,
        fill=(20, 36, 58),
        spacing=8,
    )

    cards = [
        (70, 160, 520, 420, "What", "Prepare Cat / GKP / Binomial states with high fidelity."),
        (560, 160, 1010, 420, "Why", "Bosonic states support logical qubits, sensing, and tests of quantum physics."),
        (1050, 160, 1530, 420, "How", "Closed-loop RL: propose controls -> measure -> reward -> update policy."),
    ]
    cols = [
        ((232, 242, 255), (79, 128, 214)),
        ((236, 255, 243), (56, 160, 108)),
        ((255, 245, 234), (208, 136, 45)),
    ]
    for (x1, y1, x2, y2, title, body), (fill, outline) in zip(cards, cols):
        _rr(d, (x1, y1, x2, y2), 26, fill=fill, outline=outline, width=4)
        _draw_text_in_box(
            d,
            (x1 + 24, y1 + 20, x2 - 24, y1 + 75),
            title,
            start_size=30,
            min_size=24,
            bold=True,
            fill=(30, 46, 72),
            spacing=6,
        )
        _draw_text_in_box(
            d,
            (x1 + 24, y1 + 85, x2 - 24, y2 - 20),
            body,
            start_size=23,
            min_size=16,
            fill=(33, 52, 80),
            spacing=8,
        )

    _rr(d, (70, 470, 1530, 830), 30, fill=(255, 255, 255), outline=(124, 143, 173), width=3)
    _draw_text_in_box(
        d,
        (105, 510, 360, 560),
        "Presentation roadmap",
        start_size=30,
        min_size=24,
        bold=True,
        fill=(36, 55, 85),
    )
    _draw_text_in_box(
        d,
        (105, 565, 1490, 805),
        "1) Explain RL and optimization in plain physics language.\n"
        "2) Show characteristic-function sampling logic with GIFs.\n"
        "3) Present robust-training results under dephasing noise.\n"
        "4) Make clear what works now and what still needs improvement.",
        start_size=23,
        min_size=16,
        fill=(35, 55, 85),
        spacing=11,
    )
    img.save(path)


def build_applications_image(path: Path) -> None:
    w, h = 1600, 900
    img = Image.new("RGB", (w, h), (248, 251, 255))
    d = ImageDraw.Draw(img)
    f_title = _font(50, bold=True)
    f_sub = _font(32, bold=True)
    f_txt = _font(24)

    _draw_text_in_box(
        d,
        (60, 28, 1540, 112),
        "Why Bosonic State Preparation Matters",
        start_size=50,
        min_size=36,
        bold=True,
        fill=(22, 38, 60),
    )

    cols = [
        (80, 150, 520, 820, "Logical qubits", "Hardware-efficient encoding\nError-protected bosonic code space\nBetter fault-tolerant building blocks"),
        (580, 150, 1020, 820, "Quantum sensing", "Non-classical states improve sensitivity\nPhase-space structure can amplify signal\nUseful beyond digital computing"),
        (1080, 150, 1520, 820, "Fundamental physics", "Control of non-classical states\nBenchmark for quantum dynamics\nPlatform to test control principles"),
    ]
    fills = [(232, 243, 255), (236, 255, 245), (255, 247, 236)]
    outlines = [(78, 126, 211), (56, 160, 108), (206, 136, 46)]

    for (x1, y1, x2, y2, title, body), fill, outline in zip(cols, fills, outlines):
        _rr(d, (x1, y1, x2, y2), 24, fill=fill, outline=outline, width=4)
        _draw_text_in_box(
            d,
            (x1 + 24, y1 + 22, x2 - 24, y1 + 82),
            title,
            start_size=32,
            min_size=24,
            bold=True,
            fill=(29, 46, 72),
        )
        _draw_text_in_box(
            d,
            (x1 + 24, y1 + 95, x2 - 24, y2 - 330),
            body,
            start_size=24,
            min_size=16,
            fill=(35, 55, 84),
            spacing=10,
        )

        # simple icon motifs
        cx = (x1 + x2) // 2
        if "qubits" in title.lower():
            for k in range(3):
                r = 48 + 18 * k
                d.ellipse((cx - r, y2 - 210 - r, cx + r, y2 - 210 + r), outline=(75, 126, 210), width=4)
        elif "sensing" in title.lower():
            d.polygon([(cx - 90, y2 - 180), (cx, y2 - 300), (cx + 90, y2 - 180)], outline=(58, 156, 107), fill=(208, 241, 223))
            d.line((cx, y2 - 300, cx, y2 - 120), fill=(58, 156, 107), width=5)
        else:
            d.ellipse((cx - 115, y2 - 300, cx + 115, y2 - 70), outline=(202, 128, 42), width=5)
            d.line((cx - 115, y2 - 185, cx + 115, y2 - 185), fill=(202, 128, 42), width=4)
            d.line((cx, y2 - 300, cx, y2 - 70), fill=(202, 128, 42), width=4)

    img.save(path)


def build_motivation_image(path: Path) -> None:
    w, h = 1600, 900
    img = Image.new("RGB", (w, h), (247, 250, 255))
    d = ImageDraw.Draw(img)
    f_title = _font(50, bold=True)
    f_sub = _font(30, bold=True)
    f_txt = _font(26)

    _draw_text_in_box(
        d,
        (60, 28, 1540, 110),
        "Encoding Intuition: Multi-Qubit vs Bosonic Mode",
        start_size=50,
        min_size=36,
        bold=True,
        fill=(24, 35, 52),
    )

    left = (50, 120, 760, 840)
    right = (840, 120, 1550, 840)
    _rr(d, left, 30, fill=(236, 242, 255), outline=(86, 126, 214), width=4)
    _rr(d, right, 30, fill=(236, 255, 243), outline=(44, 164, 109), width=4)

    _draw_text_in_box(
        d,
        (90, 150, 700, 215),
        "Many physical qubits",
        start_size=30,
        min_size=22,
        bold=True,
        fill=(32, 66, 140),
    )
    _draw_text_in_box(
        d,
        (900, 150, 1500, 215),
        "One bosonic mode",
        start_size=30,
        min_size=22,
        bold=True,
        fill=(19, 112, 74),
    )

    gx0, gy0 = 110, 240
    step = 72
    for r in range(7):
        for c in range(7):
            x = gx0 + c * step
            y = gy0 + r * step
            fill = (31, 41, 55)
            if (r, c) in {(1, 1), (2, 4), (5, 2), (4, 5), (3, 3)}:
                fill = (91, 132, 255)
            d.ellipse((x - 18, y - 18, x + 18, y + 18), fill=fill)

    _draw_text_in_box(
        d,
        (95, 740, 720, 810),
        "Protection often needs many coupled qubits",
        start_size=26,
        min_size=18,
        fill=(20, 35, 58),
    )

    ox0, oy0 = 930, 255
    ow, oh = 520, 430
    _rr(d, (ox0, oy0, ox0 + ow, oy0 + oh), 24, fill=(248, 255, 251), outline=(65, 157, 107), width=3)
    for k in range(8):
        y = oy0 + 40 + k * 46
        d.line((ox0 + 45, y, ox0 + ow - 45, y), fill=(142, 196, 164), width=2)
    d.ellipse((ox0 + 210, oy0 + 120, ox0 + 310, oy0 + 220), outline=(26, 130, 89), width=5)
    d.ellipse((ox0 + 205, oy0 + 265, ox0 + 315, oy0 + 375), outline=(26, 130, 89), width=5)
    d.text((ox0 + 345, oy0 + 150), "|0_L>", font=f_sub, fill=(22, 100, 70))
    d.text((ox0 + 345, oy0 + 295), "|1_L>", font=f_sub, fill=(22, 100, 70))

    _draw_text_in_box(
        d,
        (895, 740, 1520, 810),
        "Richer controllable state space in one oscillator",
        start_size=26,
        min_size=18,
        fill=(14, 84, 56),
    )

    _arrow(d, (770, 450), (830, 450), color=(72, 98, 175), width=8, head=24)
    _draw_text_in_box(
        d,
        (560, 355, 1040, 440),
        "Hardware-efficient bosonic encoding",
        start_size=27,
        min_size=18,
        bold=True,
        fill=(53, 74, 128),
        align="center",
        valign="middle",
    )

    img.save(path)


def build_timeline_image(path: Path) -> None:
    w, h = 1600, 900
    img = Image.new("RGB", (w, h), (248, 251, 255))
    d = ImageDraw.Draw(img)
    f_title = _font(44, bold=True)
    f_year = _font(28, bold=True)
    f_txt = _font(23)

    d.text((70, 30), "Positioning This Work in Literature", font=f_title, fill=(24, 39, 60))

    y = 460
    d.line((140, y, 1460, y), fill=(90, 112, 145), width=6)
    points = [
        (220, "2001", "GKP code\nproposal\n(Gottesman-Kitaev-Preskill)"),
        (620, "2022", "Model-free RL\nquantum control\n(Sivak et al., PRX)"),
        (990, "2024", "Robust trapped-ion\nbosonic state prep\n(Matsos et al., PRL)"),
        (1360, "2026", "This project:\nRL workflow for\nCat/GKP/Binomial"),
    ]
    cols = [(79, 127, 210), (124, 92, 182), (53, 162, 107), (209, 133, 43)]
    for (x, year, txt), c in zip(points, cols):
        d.ellipse((x - 15, y - 15, x + 15, y + 15), fill=c)
        d.text((x - 35, y - 70), year, font=f_year, fill=(31, 47, 72))
        _rr(d, (x - 150, y + 40, x + 150, y + 225), 18, fill=(255, 255, 255), outline=(185, 199, 220), width=2)
        d.multiline_text((x - 132, y + 70), txt, font=f_txt, fill=(33, 50, 78), spacing=7)

    d.text((100, 120), "Use references you actually discuss in the talk; avoid citation overload.", font=f_txt, fill=(33, 53, 82))
    d.text((100, 160), "Story line: prior ideas -> our adaptation -> current robust-control results.", font=f_txt, fill=(33, 53, 82))

    img.save(path)


def build_literature_comparison_image(path: Path) -> None:
    w, h = 1600, 900
    img = Image.new("RGB", (w, h), (248, 251, 255))
    d = ImageDraw.Draw(img)
    f_title = _font(42, bold=True)
    f_sub = _font(27, bold=True)
    f_txt = _font(22)

    _draw_text_in_box(
        d,
        (70, 30, 1530, 100),
        "Key References and This Project",
        start_size=42,
        min_size=32,
        bold=True,
        fill=(24, 39, 60),
    )

    cards = [
        (
            80,
            165,
            510,
            780,
            "Sivak et al. (PRX 2022)",
            "Model-free RL for quantum control.\n"
            "Core idea: optimize controls from feedback,\n"
            "without analytic gradients.",
        ),
        (
            585,
            165,
            1015,
            780,
            "Matsos et al. (PRL 2024)",
            "Robust trapped-ion bosonic-state preparation.\n"
            "Core idea: deterministic and robust\n"
            "bosonic preparation is experimentally feasible.",
        ),
        (
            1090,
            165,
            1520,
            780,
            "This Project",
            "Applies and explains RL workflow for\n"
            "Cat/GKP/Binomial trapped-ion states.\n"
            "Adds dephasing-robust static/stochastic\n"
            "training comparison.",
        ),
    ]
    fills = [(236, 244, 255), (236, 255, 244), (255, 247, 236)]
    outlines = [(79, 128, 214), (57, 160, 108), (206, 136, 44)]

    for (x1, y1, x2, y2, title, body), fill, outline in zip(cards, fills, outlines):
        _rr(d, (x1, y1, x2, y2), 24, fill=fill, outline=outline, width=4)
        _draw_text_in_box(
            d,
            (x1 + 22, y1 + 22, x2 - 22, y1 + 86),
            title,
            start_size=30,
            min_size=22,
            bold=True,
            fill=(30, 47, 73),
        )
        _draw_text_in_box(
            d,
            (x1 + 22, y1 + 92, x2 - 22, y2 - 22),
            body,
            start_size=24,
            min_size=17,
            fill=(35, 55, 84),
            spacing=9,
        )

    _arrow(d, (510, 470), (585, 470), color=(96, 118, 149), width=7)
    _arrow(d, (1015, 470), (1090, 470), color=(96, 118, 149), width=7)

    img.save(path)


def build_char_reward_mechanism_image(path: Path) -> None:
    w, h = 1600, 900
    img = Image.new("RGB", (w, h), (248, 252, 255))
    d = ImageDraw.Draw(img)
    f_title = _font(44, bold=True)
    f_sub = _font(27, bold=True)
    f_txt = _font(22)

    _draw_text_in_box(
        d,
        (70, 30, 1530, 100),
        "Characteristic-Function Reward Workflow",
        start_size=44,
        min_size=32,
        bold=True,
        fill=(24, 39, 60),
    )

    boxes = [
        (90, 200, 480, 720, "1) Choose sample points", "Select points in phase space with staged sampling."),
        (560, 200, 950, 720, "2) Measure model response", "Evaluate characteristic values at sampled points."),
        (1030, 200, 1420, 720, "3) Compute reward", "Compare with target values and aggregate weighted mismatch."),
    ]
    fills = [(232, 243, 255), (236, 255, 244), (255, 247, 236)]
    outlines = [(79, 128, 214), (57, 160, 108), (206, 136, 44)]
    for (x1, y1, x2, y2, title, body), fill, outline in zip(boxes, fills, outlines):
        _rr(d, (x1, y1, x2, y2), 24, fill=fill, outline=outline, width=4)
        _draw_text_in_box(
            d,
            (x1 + 20, y1 + 20, x2 - 20, y1 + 84),
            title,
            start_size=27,
            min_size=20,
            bold=True,
            fill=(29, 46, 72),
        )
        _draw_text_in_box(
            d,
            (x1 + 20, y1 + 90, x2 - 20, y2 - 20),
            body,
            start_size=22,
            min_size=15,
            fill=(35, 55, 84),
            spacing=8,
        )

    _arrow(d, (480, 455), (560, 455), color=(94, 117, 148), width=7)
    _arrow(d, (950, 455), (1030, 455), color=(94, 117, 148), width=7)

    d.text((95, 770), "Reward improves when sampled characteristic values align with target structure.", font=f_txt, fill=(42, 62, 93))
    d.text((95, 805), "This gives a measurement-style objective compatible with model-free RL updates.", font=f_txt, fill=(42, 62, 93))

    img.save(path)


def build_sampling_strategy_image(path: Path) -> None:
    w, h = 1600, 900
    img = Image.new("RGB", (w, h), (249, 252, 255))
    d = ImageDraw.Draw(img)
    _draw_text_in_box(
        d,
        (70, 30, 1530, 100),
        "How Characteristic-Function Points Are Chosen",
        start_size=42,
        min_size=31,
        bold=True,
        fill=(24, 39, 60),
    )

    panels = [
        (
            (80, 170, 520, 765),
            (232, 243, 255),
            (79, 128, 214),
            "Stage 1: deterministic top-k warm start",
            "Score each candidate by |target(alpha)| * |alpha|^radial_exp, then keep top-k points. "
            "This gives strong initial supervision (not random).",
        ),
        (
            (580, 170, 1020, 765),
            (236, 255, 244),
            (57, 160, 108),
            "Stage 2: radial-stratified sampling",
            "Split radius into bins. Allocate per-bin quota by total weight mass, then sample within each bin "
            "using local normalized weights.",
        ),
        (
            (1080, 170, 1520, 765),
            (255, 247, 236),
            (206, 136, 44),
            "Stage 3: denser continuation",
            "Use the same stratified rule with larger sample budget. "
            "If any quota is missing, fill from global weighted distribution.",
        ),
    ]
    for rect, fill, outline, title, body in panels:
        x1, y1, x2, y2 = rect
        _rr(d, rect, 24, fill=fill, outline=outline, width=4)
        _draw_text_in_box(
            d,
            (x1 + 20, y1 + 22, x2 - 20, y1 + 122),
            title,
            start_size=25,
            min_size=18,
            bold=True,
            fill=(29, 46, 72),
            spacing=7,
        )
        _draw_text_in_box(
            d,
            (x1 + 20, y1 + 132, x2 - 20, y1 + 305),
            body,
            start_size=20,
            min_size=14,
            fill=(35, 55, 84),
            spacing=8,
        )

        rng = np.random.default_rng(abs(hash(title)) % (2**32))
        density = 50 if "Stage 1" in title else (95 if "Stage 2" in title else 170)
        for _ in range(density):
            px = int(rng.uniform(x1 + 25, x2 - 25))
            py = int(rng.uniform(y1 + 330, y2 - 25))
            d.ellipse((px - 2, py - 2, px + 2, py + 2), fill=(70, 105, 165))

    _arrow(d, (520, 465), (580, 465), color=(96, 118, 149), width=7)
    _arrow(d, (1020, 465), (1080, 465), color=(96, 118, 149), width=7)
    _draw_text_in_box(
        d,
        (80, 790, 1520, 860),
        "Conclusion: point selection is rule-based + stochastic, not pure random selection.",
        start_size=21,
        min_size=15,
        bold=True,
        fill=(40, 60, 90),
    )
    img.save(path)


def build_rl_loop_image(path: Path) -> None:
    w, h = 1600, 900
    img = Image.new("RGB", (w, h), (250, 251, 255))
    d = ImageDraw.Draw(img)
    f_title = _font(46, bold=True)
    f_box = _font(25, bold=True)
    f_txt = _font(23)

    d.text((70, 28), "Closed-Loop RL Optimization", font=f_title, fill=(25, 37, 57))

    boxes = {
        "policy": (95, 220, 500, 380),
        "pulse": (560, 120, 1040, 280),
        "measure": (1100, 220, 1510, 380),
        "reward": (930, 505, 1420, 675),
        "update": (220, 515, 760, 705),
    }

    colors = {
        "policy": ((232, 243, 255), (79, 130, 216)),
        "pulse": ((238, 253, 245), (54, 167, 110)),
        "measure": ((255, 245, 231), (214, 137, 45)),
        "reward": ((250, 240, 255), (160, 86, 198)),
        "update": ((236, 246, 251), (65, 151, 189)),
    }

    labels = {
        "policy": "Policy outputs\ncontrol parameters",
        "pulse": "Apply pulse sequence\nto quantum system",
        "measure": "Sample characteristic\nfunction points",
        "reward": "Compute scalar reward\n(target consistency)",
        "update": "PPO updates policy\nfor next epoch",
    }

    for key, rect in boxes.items():
        fill, outline = colors[key]
        _rr(d, rect, 26, fill=fill, outline=outline, width=4)
        x1, y1, x2, y2 = rect
        bb = d.multiline_textbbox((0, 0), labels[key], font=f_box, align="center", spacing=7)
        tw = bb[2] - bb[0]
        th = bb[3] - bb[1]
        d.multiline_text(
            (x1 + (x2 - x1 - tw) // 2, y1 + (y2 - y1 - th) // 2),
            labels[key],
            font=f_box,
            fill=(25, 35, 50),
            align="center",
            spacing=7,
        )

    _arrow(d, (500, 300), (560, 205), color=(67, 120, 201), width=7)
    _arrow(d, (1040, 205), (1100, 300), color=(66, 155, 104), width=7)
    _arrow(d, (1285, 380), (1180, 505), color=(195, 122, 39), width=7)
    _arrow(d, (930, 590), (760, 620), color=(137, 83, 178), width=7)
    _arrow(d, (220, 600), (95, 340), color=(58, 141, 177), width=7)

    d.text((120, 770), "One epoch = rollout episodes -> compute reward -> multiple policy updates -> evaluate", font=f_txt, fill=(34, 56, 94))
    d.text((120, 807), "The cycle repeats until the policy reaches the target fidelity and robustness regime.", font=f_txt, fill=(34, 56, 94))

    img.save(path)


def build_epoch_timeline_image(path: Path) -> None:
    w, h = 1600, 900
    img = Image.new("RGB", (w, h), (250, 252, 255))
    d = ImageDraw.Draw(img)
    _draw_text_in_box(
        d,
        (70, 30, 1530, 105),
        "Training Units: Episode, Epoch, and Policy Update",
        start_size=44,
        min_size=32,
        bold=True,
        fill=(24, 38, 58),
    )

    defs = [
        (
            (80, 150, 500, 330),
            (232, 243, 255),
            (78, 125, 211),
            "Episode",
            "One full pulse attempt: propose controls -> evolve -> get reward.",
        ),
        (
            (560, 150, 980, 330),
            (236, 255, 244),
            (53, 160, 106),
            "Epoch",
            "A batch of episodes collected with current policy parameters.",
        ),
        (
            (1040, 150, 1520, 330),
            (255, 247, 236),
            (202, 137, 45),
            "Policy updates / epoch",
            "Reuse the same epoch data for U PPO updates before next rollout batch.",
        ),
    ]
    for rect, fill, outline, title, body in defs:
        x1, y1, x2, y2 = rect
        _rr(d, rect, 24, fill=fill, outline=outline, width=4)
        _draw_text_in_box(
            d,
            (x1 + 20, y1 + 18, x2 - 20, y1 + 76),
            title,
            start_size=28,
            min_size=20,
            bold=True,
            fill=(30, 47, 72),
        )
        _draw_text_in_box(
            d,
            (x1 + 20, y1 + 80, x2 - 20, y2 - 18),
            body,
            start_size=20,
            min_size=14,
            fill=(35, 55, 84),
            spacing=8,
        )

    _arrow(d, (500, 240), (560, 240), color=(95, 118, 155), width=6)
    _arrow(d, (980, 240), (1040, 240), color=(95, 118, 155), width=6)

    flow_boxes = [
        (110, 430, 380, 610, "Rollout episodes"),
        (430, 430, 700, 610, "Compute returns\nand advantages"),
        (750, 430, 1020, 610, "Run U PPO\nupdates"),
        (1070, 430, 1340, 610, "Evaluate and\nlog"),
    ]
    for i, (x1, y1, x2, y2, txt) in enumerate(flow_boxes):
        _rr(d, (x1, y1, x2, y2), 20, fill=(245, 248, 254), outline=(146, 167, 201), width=3)
        _draw_text_in_box(
            d,
            (x1 + 14, y1 + 14, x2 - 14, y2 - 14),
            txt,
            start_size=22,
            min_size=15,
            bold=True,
            align="center",
            valign="middle",
            fill=(34, 52, 81),
        )
        if i < len(flow_boxes) - 1:
            _arrow(d, (x2, (y1 + y2) // 2), (flow_boxes[i + 1][0], (y1 + y2) // 2), color=(96, 118, 149), width=6)

    _draw_text_in_box(
        d,
        (110, 675, 1490, 850),
        "Why this matters: U is a high-impact hyperparameter. Too small underuses expensive rollouts; "
        "too large can overfit one batch and destabilize training.",
        start_size=22,
        min_size=15,
        bold=True,
        fill=(46, 66, 97),
        spacing=8,
    )

    img.save(path)


def build_ppo_image(path: Path) -> None:
    w, h = 1600, 900
    img = Image.new("RGB", (w, h), (248, 252, 255))
    d = ImageDraw.Draw(img)
    _draw_text_in_box(
        d,
        (70, 28, 1530, 100),
        "Why PPO Is Used for This Control Problem",
        start_size=44,
        min_size=32,
        bold=True,
        fill=(26, 40, 62),
    )

    _rr(d, (90, 145, 1510, 255), 20, fill=(240, 246, 255), outline=(126, 149, 189), width=3)
    _draw_text_in_box(
        d,
        (115, 165, 1480, 235),
        "Continuous actions (phase/amplitude/duration) make actor-critic policy methods more natural than "
        "discrete-action value-iteration methods.",
        start_size=22,
        min_size=15,
        fill=(36, 55, 86),
        spacing=8,
    )

    steps = [
        (90, 300, 500, 760, "Actor", "Outputs a Gaussian action distribution for pulse parameters."),
        (595, 300, 1005, 760, "Critic", "Estimates value baseline used to compute advantage."),
        (1100, 300, 1510, 760, "Clipped update", "Improves policy while limiting how far each update can move from previous policy."),
    ]
    fills = [(232, 243, 255), (238, 255, 246), (255, 247, 236)]
    outlines = [(75, 126, 210), (54, 160, 108), (206, 136, 44)]
    for (x1, y1, x2, y2, title, body), fill, outline in zip(steps, fills, outlines):
        _rr(d, (x1, y1, x2, y2), 26, fill=fill, outline=outline, width=4)
        _draw_text_in_box(
            d,
            (x1 + 26, y1 + 22, x2 - 26, y1 + 85),
            title,
            start_size=28,
            min_size=20,
            bold=True,
            fill=(31, 48, 74),
        )
        _draw_text_in_box(
            d,
            (x1 + 26, y1 + 95, x2 - 26, y2 - 20),
            body,
            start_size=21,
            min_size=15,
            fill=(36, 56, 86),
            spacing=9,
        )

    _arrow(d, (500, 530), (595, 530), color=(92, 116, 150), width=7)
    _arrow(d, (1005, 530), (1100, 530), color=(92, 116, 150), width=7)
    _draw_text_in_box(
        d,
        (90, 790, 1510, 858),
        "Operationally: each epoch collects rollouts, then runs multiple PPO updates (policy updates per epoch).",
        start_size=21,
        min_size=15,
        bold=True,
        fill=(40, 60, 90),
    )

    img.save(path)


def build_dynamics_vs_qutip_image(path: Path) -> None:
    w, h = 1600, 900
    img = Image.new("RGB", (w, h), (248, 251, 255))
    d = ImageDraw.Draw(img)
    f_title = _font(44, bold=True)
    f_sub = _font(28, bold=True)
    f_txt = _font(23)

    d.text((70, 30), "Why Use dynamiqs + GPU (instead of QuTiP here)", font=f_title, fill=(24, 39, 60))

    _rr(d, (110, 160, 760, 760), 26, fill=(236, 244, 255), outline=(79, 128, 214), width=4)
    _rr(d, (840, 160, 1490, 760), 26, fill=(236, 255, 244), outline=(57, 160, 108), width=4)

    d.text((140, 195), "QuTiP baseline", font=f_sub, fill=(38, 61, 98))
    d.multiline_text(
        (140, 260),
        "Great for prototyping\nand small-scale simulation,\nbut throughput can be limiting\nfor many RL rollouts.",
        font=f_txt,
        fill=(38, 61, 98),
        spacing=10,
    )

    d.text((870, 195), "dynamiqs + JAX path", font=f_sub, fill=(24, 108, 74))
    d.multiline_text(
        (870, 260),
        "Vectorized simulation\nand parallel rollout execution\nfit RL training loops better.\nGPU acceleration shortens\niteration time.",
        font=f_txt,
        fill=(24, 108, 74),
        spacing=10,
    )

    _arrow(d, (760, 460), (840, 460), color=(96, 118, 149), width=8)
    _rr(d, (700, 485, 900, 535), 16, fill=(250, 252, 255), outline=(167, 184, 210), width=2)
    _draw_text_in_box(
        d,
        (708, 490, 892, 528),
        "large rollout batches",
        start_size=21,
        min_size=16,
        bold=True,
        fill=(64, 82, 112),
        align="center",
        valign="middle",
    )

    d.text((120, 810), "Practical implication: this toolchain improves parallel rollout throughput for RL training.", font=f_txt, fill=(40, 60, 90))

    img.save(path)


def build_hyperparam_image(metrics: Metrics, path: Path) -> None:
    h = metrics.hyper
    w, hh = 1600, 900
    img = Image.new("RGB", (w, hh), (249, 252, 255))
    d = ImageDraw.Draw(img)
    _draw_text_in_box(
        d,
        (70, 30, 1530, 100),
        "Why Different Targets Use Different Training Settings",
        start_size=42,
        min_size=31,
        bold=True,
        fill=(24, 39, 60),
    )

    cards = [
        (
            (80, 145, 520, 815),
            (232, 243, 255),
            (79, 128, 214),
            "Observed in runs",
            "Cat: epochs=%s, policy updates=%s\n"
            "GKP: epochs=%s, policy updates=%s\n"
            "Binomial: epochs=%s, policy updates=%s\n\n"
            "Same optimizer family, but best settings differ across targets."
            % (
                h.cat_num_epochs,
                h.cat_policy_updates,
                h.gkp_num_epochs,
                h.gkp_policy_updates,
                h.bin_num_epochs,
                h.bin_policy_updates,
            ),
        ),
        (
            (580, 145, 1020, 815),
            (236, 255, 244),
            (57, 160, 108),
            "Plausible explanation",
            "Target structures differ in complexity.\n\n"
            "GKP has periodic lattice-like structure and is highly phase sensitive.\n"
            "Binomial is sparse in Fock space and can be sensitive to sampling strategy.\n"
            "Therefore convergence speed and stability can require different epoch budgets.",
        ),
        (
            (1080, 145, 1520, 815),
            (255, 247, 236),
            (206, 136, 44),
            "What is still unknown",
            "We do not yet have a strict causal attribution for each hyperparameter difference.\n\n"
            "Current statement is empirical: these settings worked best in observed runs.\n"
            "Further ablation is needed for definitive explanation.",
        ),
    ]

    for rect, fill, outline, title, body in cards:
        x1, y1, x2, y2 = rect
        _rr(d, rect, 24, fill=fill, outline=outline, width=4)
        _draw_text_in_box(
            d,
            (x1 + 20, y1 + 20, x2 - 20, y1 + 92),
            title,
            start_size=28,
            min_size=20,
            bold=True,
            fill=(30, 47, 73),
        )
        _draw_text_in_box(
            d,
            (x1 + 20, y1 + 96, x2 - 20, y2 - 20),
            body,
            start_size=20,
            min_size=14,
            fill=(35, 55, 84),
            spacing=8,
        )

    img.save(path)


def build_robust_modes_image(path: Path) -> None:
    w, h = 1600, 900
    img = Image.new("RGB", (w, h), (248, 252, 255))
    d = ImageDraw.Draw(img)
    f_title = _font(46, bold=True)
    f_sub = _font(27, bold=True)
    f_txt = _font(22)

    _draw_text_in_box(
        d,
        (70, 30, 1530, 100),
        "Dephasing-Robust Training Modes",
        start_size=46,
        min_size=33,
        bold=True,
        fill=(24, 39, 60),
    )

    # quasi-static box
    _rr(d, (100, 160, 760, 780), 26, fill=(236, 244, 255), outline=(79, 128, 214), width=4)
    _draw_text_in_box(
        d,
        (130, 195, 730, 248),
        "Quasi-static branch",
        start_size=27,
        min_size=21,
        bold=True,
        fill=(35, 61, 98),
    )
    _draw_text_in_box(
        d,
        (130, 255, 730, 700),
        "One detuning sample per trajectory (or a detuning grid across samples).\n"
        "Configuration used here:\n"
        "- grid-based detuning samples\n"
        "- Gaussian weighting for objective aggregation\n"
        "- nominal sample included in the objective",
        start_size=22,
        min_size=15,
        fill=(35, 61, 98),
        spacing=9,
    )

    # stochastic box
    _rr(d, (840, 160, 1500, 780), 26, fill=(236, 255, 244), outline=(57, 160, 108), width=4)
    _draw_text_in_box(
        d,
        (870, 195, 1470, 248),
        "Stochastic branch",
        start_size=27,
        min_size=21,
        bold=True,
        fill=(24, 108, 74),
    )
    _draw_text_in_box(
        d,
        (870, 255, 1470, 700),
        "Detuning is resampled across time segments.\n"
        "Configuration used here:\n"
        "- gamma_dt standard-deviation mode\n"
        "- Omega and t_step kept on a consistent physical scale\n"
        "- shared trajectories across each training batch",
        start_size=22,
        min_size=15,
        fill=(24, 108, 74),
        spacing=9,
    )

    _arrow(d, (760, 470), (840, 470), color=(95, 117, 148), width=8)
    _rr(d, (730, 510, 870, 555), 14, fill=(250, 252, 255), outline=(167, 184, 210), width=2)
    _draw_text_in_box(
        d,
        (736, 516, 864, 548),
        "same RL loop",
        start_size=21,
        min_size=15,
        bold=True,
        fill=(62, 80, 109),
        align="center",
        valign="middle",
    )

    d.text((110, 810), "Both modes optimize controls under noise; they differ in noise sampling model and objective aggregation.", font=f_txt, fill=(40, 60, 90))

    img.save(path)


def build_agent_io_image(path: Path) -> None:
    w, h = 1600, 900
    img = Image.new("RGB", (w, h), (249, 252, 255))
    d = ImageDraw.Draw(img)

    _draw_text_in_box(
        d,
        (70, 30, 1530, 100),
        "What the Agent Uses and What It Updates",
        start_size=44,
        min_size=32,
        bold=True,
        fill=(24, 39, 60),
    )

    cards = [
        (
            (90, 170, 730, 360),
            (232, 243, 255),
            (79, 128, 214),
            "Context (policy input)",
            "Information provided at each segment: step/clock encoding and fixed experiment context.",
        ),
        (
            (860, 170, 1510, 360),
            (236, 255, 244),
            (57, 160, 108),
            "Action (policy output)",
            "Continuous pulse parameters: phase controls (phi_r, phi_b), optional amplitudes, and duration scale.",
        ),
        (
            (90, 430, 730, 750),
            (255, 247, 236),
            (206, 136, 44),
            "Reward",
            "After full pulse execution, characteristic-function agreement with target is converted to a scalar reward.",
        ),
        (
            (860, 430, 1510, 750),
            (245, 239, 255),
            (153, 92, 186),
            "Advantage",
            "Advantage measures how much better an action performed than critic expectation; PPO uses it to update policy.",
        ),
    ]

    for rect, fill, outline, title, body in cards:
        x1, y1, x2, y2 = rect
        _rr(d, rect, 26, fill=fill, outline=outline, width=4)
        _draw_text_in_box(
            d,
            (x1 + 22, y1 + 20, x2 - 22, y1 + 76),
            title,
            start_size=27,
            min_size=20,
            bold=True,
            fill=(29, 46, 72),
        )
        _draw_text_in_box(
            d,
            (x1 + 22, y1 + 88, x2 - 22, y2 - 20),
            body,
            start_size=22,
            min_size=15,
            fill=(35, 55, 84),
            spacing=8,
        )

    _arrow(d, (730, 265), (860, 265), color=(95, 117, 148), width=7)
    _arrow(d, (1185, 360), (1185, 430), color=(95, 117, 148), width=7)
    _arrow(d, (860, 590), (730, 590), color=(95, 117, 148), width=7)
    _arrow(d, (410, 430), (410, 360), color=(95, 117, 148), width=7)

    _draw_text_in_box(
        d,
        (90, 780, 1510, 855),
        "This is policy iteration in continuous control, not discrete value-table iteration.",
        start_size=22,
        min_size=15,
        bold=True,
        fill=(40, 60, 90),
    )

    img.save(path)


def build_output_reading_guide_image(path: Path) -> None:
    w, h = 1600, 900
    img = Image.new("RGB", (w, h), (248, 252, 255))
    d = ImageDraw.Draw(img)

    _draw_text_in_box(
        d,
        (70, 30, 1530, 100),
        "How to Read the Result Figures",
        start_size=44,
        min_size=32,
        bold=True,
        fill=(24, 39, 60),
    )

    panels = [
        (
            (90, 180, 510, 780),
            (232, 243, 255),
            (79, 128, 214),
            "1) Fidelity vs Epoch",
            "Checks whether policy updates consistently improve control quality over training.",
        ),
        (
            (590, 180, 1010, 780),
            (236, 255, 244),
            (57, 160, 108),
            "2) Characteristic Map",
            "Compares target and final-state structure in phase space.",
        ),
        (
            (1090, 180, 1510, 780),
            (255, 247, 236),
            (206, 136, 44),
            "3) Pulse Sequence",
            "Shows the learned control profile that generates the final state.",
        ),
    ]

    for rect, fill, outline, title, body in panels:
        x1, y1, x2, y2 = rect
        _rr(d, rect, 24, fill=fill, outline=outline, width=4)
        _draw_text_in_box(
            d,
            (x1 + 18, y1 + 20, x2 - 18, y1 + 92),
            title,
            start_size=27,
            min_size=20,
            bold=True,
            fill=(30, 47, 73),
        )
        _draw_text_in_box(
            d,
            (x1 + 18, y1 + 105, x2 - 18, y1 + 210),
            body,
            start_size=21,
            min_size=15,
            fill=(35, 55, 84),
            spacing=8,
        )

        # simple symbolic visuals
        if "Fidelity" in title:
            d.line((x1 + 40, y2 - 120, x2 - 40, y2 - 120), fill=(87, 114, 149), width=3)
            d.line((x1 + 40, y2 - 220, x1 + 40, y2 - 80), fill=(87, 114, 149), width=3)
            pts = [
                (x1 + 55, y2 - 95),
                (x1 + 130, y2 - 130),
                (x1 + 210, y2 - 165),
                (x1 + 300, y2 - 185),
                (x1 + 370, y2 - 205),
            ]
            d.line(pts, fill=(60, 118, 204), width=6)
            for px, py in pts:
                d.ellipse((px - 5, py - 5, px + 5, py + 5), fill=(60, 118, 204))
        elif "Characteristic" in title:
            cx = (x1 + x2) // 2
            cy = y2 - 150
            for r in [30, 60, 90]:
                d.ellipse((cx - r, cy - r, cx + r, cy + r), outline=(55, 145, 100), width=4)
            d.line((cx - 110, cy, cx + 110, cy), fill=(55, 145, 100), width=3)
            d.line((cx, cy - 110, cx, cy + 110), fill=(55, 145, 100), width=3)
        else:
            x_start = x1 + 50
            y_mid = y2 - 165
            for i in range(14):
                x = x_start + i * 23
                up = int(35 * np.sin(i * 0.65))
                d.line((x, y_mid, x, y_mid - up), fill=(204, 132, 42), width=5)
            d.line((x1 + 40, y2 - 120, x2 - 40, y2 - 120), fill=(120, 132, 150), width=2)

    _arrow(d, (510, 480), (590, 480), color=(96, 118, 149), width=7)
    _arrow(d, (1010, 480), (1090, 480), color=(96, 118, 149), width=7)

    img.save(path)


def build_next_steps_image(path: Path) -> None:
    w, h = 1600, 900
    img = Image.new("RGB", (w, h), (249, 252, 255))
    d = ImageDraw.Draw(img)

    _draw_text_in_box(
        d,
        (70, 30, 1530, 100),
        "Execution Plan Before the Talk",
        start_size=42,
        min_size=31,
        bold=True,
        fill=(24, 39, 60),
    )

    stages = [
        ((120, 190, 500, 760), "1) Stability pass", "Finalize slide wording, figure consistency, and narrative transitions."),
        ((610, 190, 990, 760), "2) Dry run", "Practice timing and verify that RL workflow explanation is clear to physics audience."),
        ((1100, 190, 1480, 760), "3) Final polish", "Incorporate advisor feedback and lock the presentation for Tuesday report."),
    ]
    fills = [(232, 243, 255), (236, 255, 244), (255, 247, 236)]
    outlines = [(79, 128, 214), (57, 160, 108), (206, 136, 44)]
    for (rect, title, body), fill, outline in zip(stages, fills, outlines):
        x1, y1, x2, y2 = rect
        _rr(d, rect, 26, fill=fill, outline=outline, width=4)
        _draw_text_in_box(
            d,
            (x1 + 20, y1 + 20, x2 - 20, y1 + 92),
            title,
            start_size=27,
            min_size=20,
            bold=True,
            fill=(29, 46, 72),
        )
        _draw_text_in_box(
            d,
            (x1 + 20, y1 + 105, x2 - 20, y2 - 20),
            body,
            start_size=22,
            min_size=15,
            fill=(35, 55, 84),
            spacing=8,
        )

    _arrow(d, (500, 475), (610, 475), color=(96, 118, 149), width=7)
    _arrow(d, (990, 475), (1100, 475), color=(96, 118, 149), width=7)
    img.save(path)


def build_conclusion_image(metrics: Metrics, path: Path) -> None:
    w, h = 1600, 900
    img = Image.new("RGB", (w, h), (248, 252, 255))
    d = ImageDraw.Draw(img)
    _draw_text_in_box(
        d,
        (70, 30, 1530, 100),
        "Project Snapshot",
        start_size=44,
        min_size=32,
        bold=True,
        fill=(24, 39, 60),
    )

    cards = [
        (
            (100, 180, 730, 420),
            "Cross-target nominal performance",
            "Cat %.4f | GKP %.4f | Binomial baseline %.4f"
            % (metrics.cat_fid, metrics.gkp_fid, metrics.bin_pre_fid),
        ),
        (
            (860, 180, 1490, 420),
            "Best robust penalties",
            "Quasi-static p=%.2f | Stochastic p=%.2f"
            % (metrics.best_static.penalty, metrics.best_stoch.penalty),
        ),
        (
            (100, 500, 730, 760),
            "Robust quality now",
            "Static score %.4f (f_nom %.4f, f_rob %.4f)"
            % (metrics.best_static.score, metrics.best_static.f_nom, metrics.best_static.f_rob),
        ),
        (
            (860, 500, 1490, 760),
            "Main open direction",
            "Raise stochastic branch robustness while preserving near-zero-detuning fidelity.",
        ),
    ]
    fills = [(232, 243, 255), (236, 255, 244), (255, 247, 236), (245, 239, 255)]
    outlines = [(79, 128, 214), (57, 160, 108), (206, 136, 44), (153, 92, 186)]
    for (rect, title, body), fill, outline in zip(cards, fills, outlines):
        x1, y1, x2, y2 = rect
        _rr(d, rect, 24, fill=fill, outline=outline, width=4)
        _draw_text_in_box(
            d,
            (x1 + 20, y1 + 18, x2 - 20, y1 + 82),
            title,
            start_size=26,
            min_size=19,
            bold=True,
            fill=(29, 46, 72),
        )
        _draw_text_in_box(
            d,
            (x1 + 20, y1 + 90, x2 - 20, y2 - 20),
            body,
            start_size=23,
            min_size=15,
            fill=(35, 55, 84),
            spacing=8,
        )
    img.save(path)


def build_reference_scope_image(path: Path) -> None:
    w, h = 1600, 900
    img = Image.new("RGB", (w, h), (249, 252, 255))
    d = ImageDraw.Draw(img)
    _draw_text_in_box(
        d,
        (70, 30, 1530, 100),
        "Reference Scope Used in This Talk",
        start_size=42,
        min_size=31,
        bold=True,
        fill=(24, 39, 60),
    )

    rows = [
        ("Model-free RL context", "Sivak et al., PRX 2022"),
        ("Bosonic trapped-ion benchmark", "Matsos et al., PRL 2024"),
        ("Bosonic code motivation", "Gottesman-Kitaev-Preskill, PRA 2001"),
        ("Project-specific results", "Current Cat/GKP/Binomial outputs and robust sweeps"),
    ]
    y = 180
    for i, (topic, cite) in enumerate(rows):
        fill = (235, 244, 255) if i % 2 == 0 else (241, 248, 255)
        _rr(d, (120, y, 1480, y + 130), 18, fill=fill, outline=(164, 184, 214), width=2)
        _draw_text_in_box(
            d,
            (150, y + 18, 620, y + 112),
            topic,
            start_size=27,
            min_size=19,
            bold=True,
            fill=(32, 49, 75),
        )
        _draw_text_in_box(
            d,
            (640, y + 18, 1450, y + 112),
            cite,
            start_size=24,
            min_size=16,
            fill=(38, 58, 89),
        )
        y += 155
    img.save(path)


def build_agenda_visual_image(path: Path) -> None:
    w, h = 1600, 900
    img = Image.new("RGB", (w, h), (248, 252, 255))
    d = ImageDraw.Draw(img)
    _draw_text_in_box(
        d,
        (70, 30, 1530, 100),
        "Report Flow Overview",
        start_size=44,
        min_size=32,
        bold=True,
        fill=(24, 39, 60),
    )
    blocks = [
        ((120, 170, 720, 360), "Motivation", "Why bosonic state preparation and why RL now."),
        ((880, 170, 1480, 360), "Method", "Closed-loop optimization, PPO updates, and sampling strategy."),
        ((120, 430, 720, 720), "Results", "Cat/GKP/Binomial outcomes with robust-training sweeps."),
        ((880, 430, 1480, 720), "Outlook", "What works now, what to improve next, and action plan."),
    ]
    fills = [(232, 243, 255), (236, 255, 244), (255, 247, 236), (245, 239, 255)]
    outlines = [(79, 128, 214), (57, 160, 108), (206, 136, 44), (153, 92, 186)]
    for (rect, title, body), fill, outline in zip(blocks, fills, outlines):
        x1, y1, x2, y2 = rect
        _rr(d, rect, 24, fill=fill, outline=outline, width=4)
        _draw_text_in_box(
            d,
            (x1 + 22, y1 + 20, x2 - 22, y1 + 84),
            title,
            start_size=29,
            min_size=21,
            bold=True,
            fill=(30, 47, 73),
        )
        _draw_text_in_box(
            d,
            (x1 + 22, y1 + 95, x2 - 22, y2 - 22),
            body,
            start_size=23,
            min_size=16,
            fill=(35, 55, 84),
            spacing=8,
        )

    _arrow(d, (720, 265), (880, 265), color=(95, 117, 148), width=7)
    _arrow(d, (1180, 360), (1180, 430), color=(95, 117, 148), width=7)
    _arrow(d, (880, 575), (720, 575), color=(95, 117, 148), width=7)
    _arrow(d, (420, 430), (420, 360), color=(95, 117, 148), width=7)
    img.save(path)


def build_scope_summary_image(path: Path) -> None:
    w, h = 1600, 900
    img = Image.new("RGB", (w, h), (248, 252, 255))
    d = ImageDraw.Draw(img)
    _draw_text_in_box(
        d,
        (70, 30, 1530, 100),
        "Project Scope Summary",
        start_size=44,
        min_size=32,
        bold=True,
        fill=(24, 39, 60),
    )

    rows = [
        ("Targets", "Cat, GKP, Binomial"),
        ("Optimization", "Model-free RL with PPO policy updates"),
        ("Reward", "Characteristic-function matching with stage-wise sampling"),
        ("Robust extension", "Quasi-static and stochastic dephasing branches"),
        ("Compute stack", "dynamiqs + GPU for parallel rollout throughput"),
    ]
    y = 170
    for i, (k, v) in enumerate(rows):
        fill = (234, 244, 255) if i % 2 == 0 else (242, 249, 255)
        _rr(d, (120, y, 1480, y + 120), 16, fill=fill, outline=(170, 190, 215), width=2)
        _draw_text_in_box(
            d,
            (150, y + 22, 500, y + 100),
            k,
            start_size=29,
            min_size=21,
            bold=True,
            fill=(34, 54, 82),
        )
        _draw_text_in_box(
            d,
            (520, y + 22, 1450, y + 100),
            v,
            start_size=25,
            min_size=17,
            fill=(39, 60, 90),
        )
        y += 135
    img.save(path)


def _plot_nonrobust_history(metrics: Metrics, path: Path) -> None:
    hist = metrics.nonrobust_history
    fig, ax = plt.subplots(figsize=(11.2, 5.4))
    if not hist:
        ax.text(0.5, 0.5, "No non-robust history found", transform=ax.transAxes, ha="center", va="center")
        ax.set_title("Binomial Non-Robust Runs")
        fig.savefig(path, dpi=220, bbox_inches="tight")
        plt.close(fig)
        return
    xs = np.arange(len(hist))
    ys = np.array([v for _, v in hist], dtype=float)
    labels = []
    for k, _ in hist:
        if k.startswith("checkpoint_"):
            labels.append("checkpoint best")
        else:
            labels.append(k.replace("202602", "02-").replace("_", " "))
    ax.plot(xs, ys, "o-", color="#2c6db7", lw=2.2)
    best_idx = int(np.argmax(ys))
    ax.scatter([best_idx], [ys[best_idx]], s=120, color="#d2862c", zorder=5, label="best non-robust")
    if metrics.checkpoint_nonrobust_fid > 0:
        ax.axhline(
            metrics.checkpoint_nonrobust_fid,
            color="#8a5fd0",
            lw=1.8,
            ls="--",
            label="checkpoint best (no-noise)",
        )
    for i, y in enumerate(ys):
        ax.text(i, y + 0.008, f"{y:.3f}", ha="center", fontsize=8)
    ax.set_ylim(max(0.0, float(np.min(ys)) - 0.05), min(1.0, float(np.max(ys)) + 0.08))
    ax.set_xticks(xs)
    ax.set_xticklabels(labels, rotation=25, ha="right", fontsize=8)
    ax.set_ylabel("Final fidelity")
    ax.set_title("Binomial No-Noise / Non-Robust History (logs + checkpoint)")
    ax.grid(alpha=0.25)
    ax.legend(loc="lower right")
    fig.tight_layout()
    fig.savefig(path, dpi=220, bbox_inches="tight")
    plt.close(fig)


def _plot_penalty_sweep(entries: List[SweepEntry], title: str, path: Path, best_penalty: float) -> None:
    if not entries:
        plt.figure(figsize=(11, 6))
        plt.title(title)
        plt.text(0.5, 0.5, "No data", ha="center", va="center")
        plt.savefig(path, dpi=200, bbox_inches="tight")
        plt.close()
        return

    xs = np.array([e.penalty for e in entries], dtype=float)
    f_nom = np.array([e.f_nom for e in entries], dtype=float)
    f_rob = np.array([e.f_rob for e in entries], dtype=float)
    score = np.array([e.score for e in entries], dtype=float)

    fig, ax = plt.subplots(figsize=(11.5, 6.5))
    ax.plot(xs, f_nom, "o-", lw=2.5, color="#2c6db7", label="Nominal fidelity")
    ax.plot(xs, f_rob, "o-", lw=2.5, color="#d2862c", label="Robust fidelity")
    ax.plot(xs, score, "o-", lw=2.5, color="#2f9b5c", label="Score")

    ax.set_xlabel("Penalty p")
    ax.set_ylabel("Metric")
    ax.set_title(title)
    ax.grid(alpha=0.25)
    ax.set_ylim(0.65, 1.01)
    ax.legend(loc="lower right")

    best_idx = int(np.argmin(np.abs(xs - best_penalty)))
    ax.axvline(xs[best_idx], color="#ba2d2d", linestyle="--", alpha=0.8)
    ax.text(xs[best_idx], 0.68, "best p=%.2f" % xs[best_idx], color="#ba2d2d")

    for x, y in zip(xs, score):
        ax.text(x, y + 0.01, "%.3f" % y, ha="center", fontsize=9)

    fig.savefig(path, dpi=220, bbox_inches="tight")
    plt.close(fig)


def _plot_eval_curve_compare(static_run: Path, stoch_run: Path, path: Path) -> None:
    fig, ax = plt.subplots(figsize=(11.5, 6.5))
    used = False

    static_csv = static_run / "eval_fidelity.csv"
    if static_csv.exists():
        df = pd.read_csv(static_csv)
        if {"epoch", "mean_fidelity"}.issubset(df.columns):
            ax.plot(df["epoch"], df["mean_fidelity"], lw=2.0, color="#2c6db7", label="Static best (p=%.2f)" % _parse_penalty_token(static_run.name.split("_")[1]))
            used = True

    stoch_csv = stoch_run / "eval_fidelity.csv"
    if stoch_csv.exists():
        df = pd.read_csv(stoch_csv)
        if {"epoch", "mean_fidelity"}.issubset(df.columns):
            ax.plot(df["epoch"], df["mean_fidelity"], lw=2.0, color="#2f9b5c", label="Stochastic best (p=%.2f)" % _parse_penalty_token(stoch_run.name.split("_")[2]))
            used = True

    if not used:
        ax.text(0.5, 0.5, "Missing eval_fidelity.csv", transform=ax.transAxes, ha="center", va="center")

    ax.set_title("Binomial Evaluation Fidelity vs Epoch")
    ax.set_xlabel("Epoch")
    ax.set_ylabel("Mean evaluation fidelity")
    ax.grid(alpha=0.25)
    ax.legend(loc="lower right")
    fig.savefig(path, dpi=220, bbox_inches="tight")
    plt.close(fig)


def _plot_state_dashboard(metrics: Metrics, path: Path) -> None:
    names = ["Cat", "GKP", "Binomial\nstatic-best", "Binomial\nstoch-best"]
    vals = [
        metrics.cat_fid,
        metrics.gkp_fid,
        metrics.best_static.f_nom,
        metrics.best_stoch.f_nom,
    ]
    cols = ["#3a71bd", "#2f9b5c", "#cf8a32", "#6b55b7"]

    fig, ax = plt.subplots(figsize=(11.5, 6.5))
    bars = ax.bar(names, vals, color=cols)
    ax.set_ylim(0.88, 1.01)
    ax.set_ylabel("Final nominal fidelity")
    ax.set_title("Current Best Fidelity Across Main Targets")
    ax.grid(axis="y", alpha=0.25)
    for b, v in zip(bars, vals):
        ax.text(b.get_x() + b.get_width() / 2.0, v + 0.002, "%.4f" % v, ha="center", va="bottom", fontsize=10)
    fig.savefig(path, dpi=220, bbox_inches="tight")
    plt.close(fig)


def _plot_static_vs_stoch(metrics: Metrics, path: Path) -> None:
    labels = [
        "Score",
        "Nominal fidelity",
        "Robust fidelity",
        "Mean gain\n(over baseline)",
        "Gain @ zero\ndetuning",
    ]
    static_vals = [
        metrics.best_static.score,
        metrics.best_static.f_nom,
        metrics.best_static.f_rob,
        metrics.static_stats.mean_gain,
        metrics.static_stats.zero_gain,
    ]
    stoch_vals = [
        metrics.best_stoch.score,
        metrics.best_stoch.f_nom,
        metrics.best_stoch.f_rob,
        metrics.stoch_stats.mean_gain,
        metrics.stoch_stats.zero_gain,
    ]

    x = np.arange(len(labels))
    width = 0.36

    fig, ax = plt.subplots(figsize=(12, 6.8))
    ax.bar(x - width / 2, static_vals, width, label="Static best (p=%.2f)" % metrics.best_static.penalty, color="#2c6db7")
    ax.bar(x + width / 2, stoch_vals, width, label="Stochastic best (p=%.2f)" % metrics.best_stoch.penalty, color="#2f9b5c")

    ax.set_xticks(x)
    ax.set_xticklabels(labels)
    ax.set_ylim(-0.08, 1.02)
    ax.grid(axis="y", alpha=0.25)
    ax.set_title("Best Static vs Best Stochastic Robust Training")
    ax.legend(loc="upper right")

    for i, (sv, tv) in enumerate(zip(static_vals, stoch_vals)):
        ax.text(i - width / 2, sv + 0.015, "%.3f" % sv, ha="center", fontsize=9)
        ax.text(i + width / 2, tv + 0.015, "%.3f" % tv, ha="center", fontsize=9)

    fig.savefig(path, dpi=220, bbox_inches="tight")
    plt.close(fig)


def _plot_penalty_mode_comparison(metrics: Metrics, path: Path) -> None:
    fig, axes = plt.subplots(1, 3, figsize=(15.2, 4.8), sharex=False)

    s = sorted(metrics.static_entries, key=lambda e: e.penalty)
    t = sorted(metrics.stoch_entries, key=lambda e: e.penalty)
    xs = np.array([e.penalty for e in s], dtype=float) if s else np.array([])
    xt = np.array([e.penalty for e in t], dtype=float) if t else np.array([])

    configs = [
        ("Score", lambda e: e.score, (0.72, 1.01)),
        ("Nominal fidelity", lambda e: e.f_nom, (0.90, 1.01)),
        ("Robust fidelity", lambda e: e.f_rob, (0.75, 1.01)),
    ]

    for ax, (label, fn, ylim) in zip(axes, configs):
        if s:
            ys = np.array([fn(e) for e in s], dtype=float)
            ax.plot(xs, ys, "o-", color="#2c6db7", lw=2.2, label="Quasi-static")
            for x, y in zip(xs, ys):
                ax.text(x, y + 0.006, f"{y:.3f}", ha="center", fontsize=8, color="#1f4f87")
        if t:
            yt = np.array([fn(e) for e in t], dtype=float)
            ax.plot(xt, yt, "s-", color="#2f9b5c", lw=2.2, label="Stochastic")
            for x, y in zip(xt, yt):
                ax.text(x, y - 0.018, f"{y:.3f}", ha="center", fontsize=8, color="#1f6f43")

        ax.set_title(label)
        ax.set_xlabel("Penalty p")
        ax.set_ylim(*ylim)
        ax.grid(alpha=0.25)

    axes[0].set_ylabel("Metric value")
    axes[0].legend(loc="lower left")
    fig.suptitle("Penalty Sweep Comparison Across Modes", fontsize=14, fontweight="bold")
    fig.tight_layout(rect=[0, 0, 1, 0.95])
    fig.savefig(path, dpi=220, bbox_inches="tight")
    plt.close(fig)


def _plot_binomial_progression(metrics: Metrics, path: Path) -> None:
    labels = ["Non-robust\nbest", "Static robust\nbest", "Stochastic robust\nbest"]
    fids = [metrics.best_nonrobust_fid, metrics.best_static.f_nom, metrics.best_stoch.f_nom]
    x = np.arange(len(labels))

    fig, ax = plt.subplots(figsize=(10.5, 6.2))
    ax.plot(x, fids, "o-", color="#2c6db7", lw=2.6, label="Final fidelity")
    ax.scatter([1, 2], [metrics.best_static.f_rob, metrics.best_stoch.f_rob], color="#d2862c", s=90, zorder=5, label="Robust fidelity (robust runs)")

    ax.set_xticks(x)
    ax.set_xticklabels(labels)
    ax.set_ylim(0.72, 1.01)
    ax.set_ylabel("Fidelity")
    ax.set_title("Binomial Progression: Non-Robust to Robust Branches")
    ax.grid(axis="y", alpha=0.25)
    ax.legend(loc="lower right")

    for xi, v in enumerate(fids):
        ax.text(xi, v + 0.007, f"{v:.3f}", ha="center", fontsize=9)
    ax.text(1, metrics.best_static.f_rob + 0.007, f"{metrics.best_static.f_rob:.3f}", ha="center", fontsize=9, color="#a86415")
    ax.text(2, metrics.best_stoch.f_rob + 0.007, f"{metrics.best_stoch.f_rob:.3f}", ha="center", fontsize=9, color="#a86415")

    fig.savefig(path, dpi=220, bbox_inches="tight")
    plt.close(fig)


def _plot_constant_amp_evidence(metrics: Metrics, path: Path) -> None:
    labels = [
        "Checkpoint best\n(non-robust)",
        "Strict const-amp\nnon-robust",
        "Strict const-amp\nrobust",
    ]
    vals = [
        metrics.checkpoint_nonrobust_fid,
        metrics.constamp_nonrobust.fidelity if metrics.constamp_nonrobust is not None else 0.0,
        metrics.constamp_robust.fidelity if metrics.constamp_robust is not None else 0.0,
    ]
    cols = ["#2c6db7", "#5b87c8", "#d2862c"]
    x = np.arange(len(labels))

    fig, ax = plt.subplots(figsize=(10.2, 5.6))
    bars = ax.bar(x, vals, color=cols)
    ax.set_xticks(x)
    ax.set_xticklabels(labels)
    ax.set_ylim(0.65, 1.01)
    ax.set_ylabel("Final fidelity")
    ax.set_title("Binomial Fidelity: Baseline vs Strict Constant-Amplitude Evidence")
    ax.grid(axis="y", alpha=0.25)
    for b, v in zip(bars, vals):
        ax.text(b.get_x() + b.get_width() / 2, v + 0.008, f"{v:.6f}", ha="center", fontsize=9)

    ax.text(
        0.01,
        0.02,
        "Checkpoint pulse amp std: amp_r=%.4f, amp_b=%.4f (not strict constant amplitude)"
        % (metrics.checkpoint_amp_r_std, metrics.checkpoint_amp_b_std),
        transform=ax.transAxes,
        fontsize=9,
        color="#3d4f70",
    )
    fig.savefig(path, dpi=220, bbox_inches="tight")
    plt.close(fig)


def build_penalty_formula_image(path: Path) -> None:
    w, h = 1600, 900
    img = Image.new("RGB", (w, h), (249, 252, 255))
    d = ImageDraw.Draw(img)
    _draw_text_in_box(
        d,
        (70, 30, 1530, 100),
        "Penalty Meaning in Robust Training",
        start_size=44,
        min_size=32,
        bold=True,
        fill=(24, 39, 60),
    )
    _rr(d, (90, 140, 1510, 320), 20, fill=(238, 245, 255), outline=(116, 141, 188), width=3)
    _draw_text_in_box(
        d,
        (120, 180, 1480, 275),
        "score = f_rob - p * max(0, floor - f_nom),   floor = 0.985",
        start_size=36,
        min_size=24,
        bold=True,
        align="center",
        valign="middle",
        fill=(28, 44, 70),
    )

    cards = [
        (
            (100, 380, 500, 790),
            (232, 243, 255),
            (79, 128, 214),
            "f_nom",
            "Nominal fidelity at zero detuning.\nIf this drops below floor, penalty activates.",
        ),
        (
            (600, 380, 1000, 790),
            (236, 255, 244),
            (57, 160, 108),
            "f_rob",
            "Robust fidelity under detuning samples with configured weighting.",
        ),
        (
            (1100, 380, 1500, 790),
            (255, 247, 236),
            (206, 136, 44),
            "p (penalty weight)",
            "Controls how strongly we penalize dropping below the nominal-fidelity floor.",
        ),
    ]
    for rect, fill, outline, title, body in cards:
        x1, y1, x2, y2 = rect
        _rr(d, rect, 24, fill=fill, outline=outline, width=4)
        _draw_text_in_box(
            d,
            (x1 + 20, y1 + 20, x2 - 20, y1 + 90),
            title,
            start_size=28,
            min_size=20,
            bold=True,
            fill=(30, 47, 73),
        )
        _draw_text_in_box(
            d,
            (x1 + 20, y1 + 100, x2 - 20, y2 - 20),
            body,
            start_size=20,
            min_size=14,
            fill=(35, 55, 84),
            spacing=8,
        )
    img.save(path)


def _plot_eval_curve_single(run_dir: Path, path: Path, title: str) -> None:
    fig, ax = plt.subplots(figsize=(9.2, 5.4))
    csv_path = run_dir / "eval_fidelity.csv"
    if csv_path.exists():
        df = pd.read_csv(csv_path)
        if {"epoch", "mean_fidelity"}.issubset(df.columns):
            ax.plot(df["epoch"], df["mean_fidelity"], lw=2.2, color="#2c6db7")
            ax.set_title(title)
            ax.set_xlabel("Epoch")
            ax.set_ylabel("Mean evaluation fidelity")
            ax.grid(alpha=0.25)
            fig.savefig(path, dpi=220, bbox_inches="tight")
            plt.close(fig)
            return
    ax.text(0.5, 0.5, "Missing eval_fidelity.csv", transform=ax.transAxes, ha="center", va="center")
    ax.set_title(title)
    fig.savefig(path, dpi=220, bbox_inches="tight")
    plt.close(fig)


def generate_assets(metrics: Metrics) -> None:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)

    build_applications_image(ASSET_DIR / "applications_motivation.png")
    build_motivation_image(ASSET_DIR / "motivation_encoding.png")
    build_rl_loop_image(ASSET_DIR / "rl_loop.png")
    build_agent_io_image(ASSET_DIR / "agent_io.png")
    build_epoch_timeline_image(ASSET_DIR / "epoch_timeline.png")
    build_ppo_image(ASSET_DIR / "ppo_method.png")
    build_char_reward_mechanism_image(ASSET_DIR / "char_reward_mechanism.png")
    build_sampling_strategy_image(ASSET_DIR / "sampling_strategy.png")
    build_output_reading_guide_image(ASSET_DIR / "output_reading_guide.png")
    build_conclusion_image(metrics, ASSET_DIR / "conclusion_summary.png")
    build_hyperparam_image(metrics, ASSET_DIR / "hyperparam_compare.png")
    build_dynamics_vs_qutip_image(ASSET_DIR / "dynamics_vs_qutip.png")
    build_robust_modes_image(ASSET_DIR / "robust_modes.png")
    build_penalty_formula_image(ASSET_DIR / "penalty_formula.png")

    _plot_penalty_sweep(
        metrics.static_entries,
        "Quasi-static penalty sweep (fine scan)",
        ASSET_DIR / "static_penalty_sweep.png",
        metrics.best_static.penalty,
    )
    _plot_penalty_sweep(
        metrics.stoch_entries,
        "Stochastic penalty sweep (combined coarse+refine)",
        ASSET_DIR / "stoch_penalty_sweep.png",
        metrics.best_stoch.penalty,
    )
    _plot_eval_curve_compare(metrics.best_static.run_dir, metrics.best_stoch.run_dir, ASSET_DIR / "binomial_eval_compare.png")
    _plot_eval_curve_single(metrics.bin_pre_run_dir, ASSET_DIR / "binomial_pre_eval_curve.png", "Binomial baseline checkpoint: eval fidelity")
    _plot_nonrobust_history(metrics, ASSET_DIR / "binomial_nonrobust_history.png")
    _plot_state_dashboard(metrics, ASSET_DIR / "state_dashboard.png")
    _plot_static_vs_stoch(metrics, ASSET_DIR / "static_vs_stoch.png")
    _plot_penalty_mode_comparison(metrics, ASSET_DIR / "penalty_mode_compare.png")
    _plot_binomial_progression(metrics, ASSET_DIR / "binomial_progression.png")
    _plot_constant_amp_evidence(metrics, ASSET_DIR / "binomial_constant_amp_evidence.png")

    missing = ASSET_DIR / "_missing.png"
    if not missing.exists():
        _make_missing_image(missing)


def _set_title(slide, title: str, subtitle: str = "") -> None:
    tx = slide.shapes.add_textbox(Inches(0.45), Inches(0.10), Inches(12.3), Inches(0.95))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 33, 57)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(75, 90, 114)


def _add_bullets(slide, x: float, y: float, w: float, h: float, lines: List[str], size: int = 19) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, t in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = t
        p.font.size = Pt(size)
        p.level = 0
        p.font.color.rgb = RGBColor(28, 42, 66)
        p.space_after = Pt(6)


def _add_footer(slide, page_num: int, ref_text: str = "") -> None:
    box_l = slide.shapes.add_textbox(Inches(0.35), Inches(7.10), Inches(10.8), Inches(0.26))
    tf_l = box_l.text_frame
    tf_l.clear()
    p = tf_l.paragraphs[0]
    p.text = ref_text if ref_text else "Sydney Uni weekly report | trapped-ion RL control"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(95, 110, 135)

    box_r = slide.shapes.add_textbox(Inches(12.15), Inches(7.08), Inches(0.9), Inches(0.28))
    tf_r = box_r.text_frame
    tf_r.clear()
    p2 = tf_r.paragraphs[0]
    p2.text = str(page_num)
    p2.font.size = Pt(11)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(95, 110, 135)


def _add_inline_citation(slide, text: str, x: float = 9.1, y: float = 6.78, w: float = 3.6, h: float = 0.26) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.color.rgb = RGBColor(100, 114, 136)


def _add_picture(slide, image_path: Path, x: float, y: float, w: Optional[float] = None, h: Optional[float] = None) -> None:
    fallback = ASSET_DIR / "_missing.png"
    image_path = _safe_image(image_path, fallback)
    if w is not None and h is not None:
        slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    elif w is not None:
        slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), width=Inches(w))
    elif h is not None:
        slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), height=Inches(h))
    else:
        slide.shapes.add_picture(str(image_path), Inches(x), Inches(y))


def build_presentation(metrics: Metrics) -> int:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def new_slide(title: str, subtitle: str = "", footer_ref: str = ""):
        s = prs.slides.add_slide(blank)
        _set_title(s, title, subtitle)
        _add_footer(s, len(prs.slides), footer_ref)
        return s

    cat_char = CAT_OUT / "char_target_vs_final.png"
    cat_train = CAT_OUT / "training_curve.png"
    cat_pulse = CAT_OUT / "pulse_sequences.png"
    cat_gif = CAT_OUT / "characteristic_points_evolution.gif"

    gkp_char = GKP_OUT / "char_target_vs_final.png"
    gkp_train = GKP_OUT / "training_curve.png"
    gkp_pulse = GKP_OUT / "pulse_sequences.png"
    gkp_gif = GKP_OUT / "characteristic_points_evolution.gif"

    # Use best static for robust binomial visualization, and global GIF from latest outputs.
    bin_pre_dir = metrics.bin_pre_run_dir
    bin_char_pre = bin_pre_dir / "char_target_vs_final.png"
    bin_pre_dephase = bin_pre_dir / "dephasing_compare.png"
    bin_pre_eval_curve = ASSET_DIR / "binomial_pre_eval_curve.png"

    bin_char_static = metrics.best_static.run_dir / "char_target_vs_final.png"
    bin_dephase_static = metrics.best_static.run_dir / "dephasing_compare.png"
    bin_char_stoch = metrics.best_stoch.run_dir / "char_target_vs_final.png"
    bin_dephase_stoch = metrics.best_stoch.run_dir / "dephasing_compare.png"
    bin_gif = BIN_OUT / "characteristic_points_evolution.gif"

    floor = 0.985
    penalty_example: Optional[SweepEntry] = None
    for e in sorted(metrics.static_entries + metrics.stoch_entries, key=lambda x: x.penalty):
        if e.f_nom < floor and e.penalty > 0:
            penalty_example = e
            break
    if penalty_example is None:
        penalty_example = metrics.best_static
    example_pen = penalty_example.penalty * max(0.0, floor - penalty_example.f_nom)

    const_nonrob_tag = metrics.constamp_nonrobust.tag if metrics.constamp_nonrobust is not None else "not_found"
    const_nonrob_f = metrics.constamp_nonrobust.fidelity if metrics.constamp_nonrobust is not None else 0.0
    const_nonrob_path = (
        str(metrics.constamp_nonrobust.path.relative_to(ROOT))
        if metrics.constamp_nonrobust is not None
        else "examples/trapped_ion_binomial/outputs/logs/<missing>"
    )
    const_rob_tag = metrics.constamp_robust.tag if metrics.constamp_robust is not None else "not_found"
    const_rob_f = metrics.constamp_robust.fidelity if metrics.constamp_robust is not None else 0.0
    const_rob_path = (
        str(metrics.constamp_robust.path.relative_to(ROOT))
        if metrics.constamp_robust is not None
        else "examples/trapped_ion_binomial/outputs/logs/<missing>"
    )

    # 1
    s = new_slide(
        "Trapped-Ion Bosonic State Preparation with Model-Free Reinforcement Learning",
        "Sydney Nano weekly update | 3 March 2026",
        "Project title",
    )
    _add_bullets(
        s,
        0.75,
        1.4,
        8.0,
        3.7,
        [
            "Core question: can reinforcement learning discover pulse sequences that prepare high-fidelity bosonic states?",
            "Targets in this project: Cat, GKP, and Binomial states.",
            "Extension in this report: dephasing-robust training under quasi-static and stochastic noise.",
        ],
        size=20,
    )
    _add_picture(s, ASSET_DIR / "state_dashboard.png", 8.35, 1.2, w=4.6, h=4.95)
    _add_bullets(
        s,
        0.75,
        5.45,
        12.0,
        0.9,
        ["Talk flow: motivation -> method -> RL loop -> results -> robust extension -> conclusion."],
        size=16,
    )

    # 2
    s = new_slide("Why This Problem Matters", footer_ref="Motivation")
    _add_bullets(
        s,
        0.75,
        1.05,
        5.6,
        4.9,
        [
            "Bosonic states are useful for logical qubit encodings in oscillator modes.",
            "They also matter for quantum sensing and non-classical-state engineering.",
            "Good pulse design is hard because the control landscape is high-dimensional.",
            "An automatic optimizer can accelerate discovery and reduce manual tuning cycles.",
        ],
        size=18,
    )
    _add_picture(s, ASSET_DIR / "applications_motivation.png", 6.35, 1.0, w=6.45)
    _add_inline_citation(s, "Refs: Gottesman-Kitaev-Preskill 2001; Matsos et al., PRL 2024", x=8.2, y=6.78, w=4.8)

    # 3
    s = new_slide("Why Model-Free RL Here?", footer_ref="Method motivation")
    _add_bullets(
        s,
        0.75,
        1.1,
        5.8,
        4.9,
        [
            "The optimization objective is obtained from measurement-like outputs, not analytic gradients.",
            "Model-free RL can improve control policies directly from rollout feedback.",
            "This makes it suitable for closed-loop settings where simulator/experiment returns rewards.",
            "In this project, RL acts as the pulse-search engine over continuous control parameters.",
        ],
        size=18,
    )
    _add_picture(s, ASSET_DIR / "motivation_encoding.png", 6.35, 1.0, w=6.45)
    _add_inline_citation(s, "Ref: Sivak et al., PRX 2022", x=9.2, y=6.78, w=3.6)

    # 4
    s = new_slide("Closed-Loop RL Control Workflow", footer_ref="Optimization loop")
    _add_bullets(
        s,
        0.75,
        1.1,
        5.5,
        4.5,
        [
            "Policy proposes segmented pulse parameters.",
            "Quantum simulation returns trajectory outcomes and reward signal.",
            "Agent updates parameters and repeats over epochs.",
            "The loop continues until fidelity and robustness targets are reached.",
        ],
        size=18,
    )
    _add_picture(s, ASSET_DIR / "rl_loop.png", 6.2, 1.0, w=6.6)

    # 5
    s = new_slide("Key RL Terms Used in This Talk", footer_ref="Term definitions")
    _add_bullets(
        s,
        0.75,
        1.1,
        5.5,
        4.6,
        [
            "Context: per-step information fed to policy (clock/step encoding and fixed settings).",
            "Action: continuous pulse parameters output by policy.",
            "Reward: scalar signal from characteristic-function agreement with target.",
            "Advantage: how much better an action was than critic expectation.",
        ],
        size=17,
    )
    _add_picture(s, ASSET_DIR / "agent_io.png", 6.1, 1.0, w=6.7)

    # 6
    s = new_slide("Episode, Epoch, and Policy Updates", footer_ref="Training units")
    _add_picture(s, ASSET_DIR / "epoch_timeline.png", 0.8, 0.95, w=11.8)

    # 7
    s = new_slide("Why PPO Is the Optimizer", footer_ref="Policy optimization choice")
    _add_bullets(
        s,
        0.75,
        1.05,
        5.3,
        4.7,
        [
            "PPO is on-policy actor-critic with clipped policy updates.",
            "It is stable for iterative improvement in continuous action spaces.",
            "This task uses continuous phase/amplitude/duration controls, so this is a natural fit.",
            "This is policy iteration, not value iteration over discrete actions.",
        ],
        size=17,
    )
    _add_picture(s, ASSET_DIR / "ppo_method.png", 6.05, 1.0, w=6.75)
    _add_inline_citation(s, "Ref: Sivak et al., PRX 2022", x=9.2, y=6.78, w=3.6)

    # 8
    s = new_slide("Why Hyperparameters Differ Across Target States", footer_ref="Evidence, hypothesis, unknown")
    _add_bullets(
        s,
        0.75,
        1.05,
        5.2,
        4.7,
        [
            "Observed fact: best settings differ for Cat, GKP, and Binomial training.",
            "Likely reason: target structures create different optimization landscapes.",
            "GKP appears more sensitive to phase alignment and needs larger training budget.",
            "Strict causal attribution is still open; current statement is empirical.",
        ],
        size=17,
    )
    _add_picture(s, ASSET_DIR / "hyperparam_compare.png", 6.1, 1.0, w=6.7)

    # 9
    s = new_slide("Why dynamiqs + GPU for This Pipeline", footer_ref="Simulation throughput choice")
    _add_bullets(
        s,
        0.75,
        1.1,
        5.3,
        4.4,
        [
            "RL repeatedly evaluates many trajectories, so rollout throughput is critical.",
            "dynamiqs/JAX supports vectorized batches and GPU acceleration.",
            "Compared with small-scale workflow, this reduces wall-clock time per training cycle.",
            "This is a training-enabler, not just an implementation detail.",
        ],
        size=17,
    )
    _add_picture(s, ASSET_DIR / "dynamics_vs_qutip.png", 6.1, 1.0, w=6.7)

    # 10
    s = new_slide("Characteristic-Function Reward", footer_ref="Reward design")
    _add_bullets(
        s,
        0.75,
        1.05,
        5.2,
        4.6,
        [
            "Reward compares sampled characteristic values between target and generated states.",
            "This directly links optimization to physically meaningful observables.",
            "The same representation is used to interpret final quality in result plots.",
            "Training and interpretation therefore use a consistent language.",
        ],
        size=17,
    )
    _add_picture(s, ASSET_DIR / "char_reward_mechanism.png", 6.1, 1.0, w=6.7)

    # 11
    s = new_slide("How Sample Points Expand During Training", footer_ref="Stage schedule + selection mechanism")
    _add_bullets(
        s,
        0.75,
        1.05,
        5.2,
        4.7,
        [
            "Stage 1 uses top-k points by weighted target magnitude for strong initial supervision.",
            "Later stages use radial bins, quota allocation, and weighted random draws in each bin.",
            "The process combines deterministic structure and stochastic sampling.",
            "This is why point evolution is guided, not purely random.",
        ],
        size=17,
    )
    _add_picture(s, ASSET_DIR / "sampling_strategy.png", 6.1, 1.0, w=6.7)

    # 12
    s = new_slide("How to Read the Result Figures", footer_ref="Result interpretation guide")
    _add_bullets(
        s,
        0.75,
        1.05,
        5.2,
        4.6,
        [
            "Fidelity curve: whether training improves over epochs.",
            "Characteristic map: whether final state reproduces target structure.",
            "Pulse plot: what control waveform was learned.",
            "Dephasing curve: robustness tradeoff versus detuning.",
        ],
        size=17,
    )
    _add_picture(s, ASSET_DIR / "output_reading_guide.png", 6.1, 1.0, w=6.7)

    # 13
    s = new_slide("Cat State Result", footer_ref="Cat target")
    _add_bullets(s, 0.75, 1.03, 12.0, 0.9, ["Final Cat fidelity: %.6f" % metrics.cat_fid], size=22)
    _add_picture(s, cat_char, 0.7, 1.75, w=4.1)
    _add_picture(s, cat_train, 4.95, 1.75, w=4.1)
    _add_picture(s, cat_pulse, 9.2, 1.75, w=4.1)
    _add_bullets(s, 0.75, 6.35, 12.0, 0.7, ["Target/final map, training trend, and learned pulse are shown side-by-side."], size=15)

    # 14
    s = new_slide("GKP State Result", footer_ref="GKP target")
    _add_bullets(s, 0.75, 1.03, 12.0, 0.9, ["Final GKP fidelity: %.6f" % metrics.gkp_fid], size=22)
    _add_picture(s, gkp_char, 0.7, 1.75, w=4.1)
    _add_picture(s, gkp_train, 4.95, 1.75, w=4.1)
    _add_picture(s, gkp_pulse, 9.2, 1.75, w=4.1)
    _add_bullets(
        s,
        0.75,
        6.35,
        12.0,
        0.7,
        ["GKP needs alignment of many periodic peaks and is more sensitive to phase errors, so convergence is harder."],
        size=15,
    )

    # 15
    s = new_slide("Binomial Baseline and Constant-Amplitude Evidence", footer_ref="Baseline clarity before robust extension")
    _add_bullets(
        s,
        0.75,
        1.03,
        5.2,
        5.0,
        [
            "Checkpoint baseline (before robust extension): %.6f." % metrics.checkpoint_nonrobust_fid,
            "This baseline pulse is not strict constant-amplitude (checkpoint amp std: r=%.4f, b=%.4f)."
            % (metrics.checkpoint_amp_r_std, metrics.checkpoint_amp_b_std),
            "Strict constant-amplitude non-robust best: %s -> %.6f." % (const_nonrob_tag, const_nonrob_f),
            "Strict constant-amplitude robust best: %s -> %.6f." % (const_rob_tag, const_rob_f),
            "Difference mainly comes from whether final refinement allows amplitude optimization.",
        ],
        size=15,
    )
    _add_picture(s, ASSET_DIR / "binomial_constant_amp_evidence.png", 5.95, 1.05, w=6.85)
    _add_bullets(
        s,
        0.75,
        6.12,
        12.0,
        0.95,
        [
            "Strict const-amp logs: %s | %s"
            % (const_nonrob_path, const_rob_path),
        ],
        size=12,
    )

    # 16
    s = new_slide("Characteristic-Point Evolution (GIFs)", footer_ref="Cat / GKP / Binomial point evolution")
    _add_bullets(
        s,
        0.75,
        1.0,
        12.0,
        0.8,
        ["Point distributions move from broad coverage to more informative regions as training progresses."],
        size=19,
    )
    _add_picture(s, cat_gif, 0.55, 1.78, w=4.15)
    _add_picture(s, gkp_gif, 4.65, 1.78, w=4.15)
    _add_picture(s, bin_gif, 8.75, 1.78, w=4.15)
    _add_bullets(s, 0.68, 6.38, 12.0, 0.6, ["Cat                     GKP                     Binomial"], size=16)

    # 17
    s = new_slide("Dephasing-Robust Training Configuration", footer_ref="Quasi-static and stochastic branches")
    _add_bullets(
        s,
        0.75,
        1.02,
        5.3,
        4.7,
        [
            "Quasi-static branch: one detuning sample per trajectory (or grid across samples).",
            "Stochastic branch: segment-wise random detuning with gamma_dt scaling.",
            "Teacher-aligned physical scale: Omega_r = Omega_b = 2pi x 2000 Hz, with T_STEP adjusted consistently.",
            "Both branches use the same PPO loop and same reward framework.",
        ],
        size=16,
    )
    _add_picture(s, ASSET_DIR / "robust_modes.png", 6.05, 1.0, w=6.75)
    _add_inline_citation(s, "Ref: Matsos et al., PRL 2024", x=9.05, y=6.78, w=3.8)

    # 18
    s = new_slide("What the Penalty Means in Robust Training", footer_ref="Robust score definition")
    _add_bullets(
        s,
        0.75,
        1.02,
        5.1,
        4.8,
        [
            "Implemented score: score = f_rob - p * max(0, floor - f_nom), with floor = 0.985.",
            "If nominal fidelity stays above floor, penalty is zero.",
            "If nominal fidelity drops below floor, score is reduced to prevent nominal collapse.",
            "Example from sweep: p=%.2f, f_nom=%.4f gives penalty term %.6f."
            % (penalty_example.penalty, penalty_example.f_nom, example_pen),
        ],
        size=16,
    )
    _add_picture(s, ASSET_DIR / "penalty_formula.png", 6.15, 1.0, w=6.65)

    # 19
    s = new_slide("Quasi-Static Penalty Sweep", footer_ref="Fine scan in static branch")
    _add_bullets(
        s,
        0.75,
        1.0,
        5.2,
        4.6,
        [
            "Best static point: p=%.2f, score=%.6f." % (metrics.best_static.penalty, metrics.best_static.score),
            "At this point: f_nom=%.6f, f_rob=%.6f." % (metrics.best_static.f_nom, metrics.best_static.f_rob),
            "Mean robust gain over baseline in detuning sweep: %.6f." % metrics.static_stats.mean_gain,
            "This branch currently gives the strongest robust average in our results.",
        ],
        size=16,
    )
    _add_picture(s, ASSET_DIR / "static_penalty_sweep.png", 6.05, 1.0, w=6.75)

    # 20
    s = new_slide("Stochastic Penalty Sweep", footer_ref="Combined stochastic scans")
    _add_bullets(
        s,
        0.75,
        1.0,
        5.2,
        4.6,
        [
            "Best stochastic point: p=%.2f, score=%.6f." % (metrics.best_stoch.penalty, metrics.best_stoch.score),
            "At this point: f_nom=%.6f, f_rob=%.6f." % (metrics.best_stoch.f_nom, metrics.best_stoch.f_rob),
            "Mean robust gain over baseline in detuning sweep: %.6f." % metrics.stoch_stats.mean_gain,
            "This branch is valid but currently below static best in this run set.",
        ],
        size=16,
    )
    _add_picture(s, ASSET_DIR / "stoch_penalty_sweep.png", 6.05, 1.0, w=6.75)

    # 21
    s = new_slide("Static vs Stochastic: Combined View", footer_ref="Mode-level comparison")
    _add_bullets(
        s,
        0.75,
        1.0,
        5.3,
        4.4,
        [
            "Both branches improve dephasing robustness compared with baseline references.",
            "Current best score remains static p=%.2f versus stochastic p=%.2f."
            % (metrics.best_static.penalty, metrics.best_stoch.penalty),
            "This comparison motivates targeted stochastic follow-up rather than abandoning stochastic mode.",
            "Practical next step is local retuning around the best stochastic region.",
        ],
        size=16,
    )
    _add_picture(s, ASSET_DIR / "penalty_mode_compare.png", 5.9, 1.0, w=6.9)
    _add_picture(s, ASSET_DIR / "binomial_eval_compare.png", 0.95, 4.85, w=4.35)

    # 22
    s = new_slide("Best Static and Best Stochastic Runs", footer_ref="Best-run detuning curves")
    _add_bullets(
        s,
        0.75,
        1.0,
        12.0,
        0.8,
        [
            "Left panel: best static run (p=%.2f). Right panel: best stochastic run (p=%.2f)."
            % (metrics.best_static.penalty, metrics.best_stoch.penalty)
        ],
        size=18,
    )
    _add_picture(s, bin_dephase_static, 0.7, 1.8, w=6.2)
    _add_picture(s, bin_dephase_stoch, 6.45, 1.8, w=6.2)
    _add_bullets(
        s,
        0.75,
        6.25,
        12.0,
        0.7,
        [
            "Static currently has better robustness average; stochastic retains a realistic random-noise formulation and remains under active tuning.",
        ],
        size=15,
    )

    # 23
    s = new_slide("Across Targets: Fidelity Snapshot", footer_ref="Cat, GKP, Binomial status")
    _add_bullets(
        s,
        0.75,
        1.0,
        12.0,
        0.8,
        [
            "Cat and GKP nominal fidelities are strong; binomial robust optimization is the current focus.",
        ],
        size=18,
    )
    _add_picture(s, ASSET_DIR / "state_dashboard.png", 1.95, 1.45, w=9.4)

    # 24
    s = new_slide("Conclusion", footer_ref="Final conclusion")
    _add_bullets(
        s,
        0.8,
        1.15,
        7.2,
        4.8,
        [
            "Model-free RL provides a clear closed-loop workflow for trapped-ion bosonic state preparation.",
            "The method works across Cat/GKP/Binomial and supports dephasing-robust extensions.",
            "Current robust best points: static p=%.2f, stochastic p=%.2f."
            % (metrics.best_static.penalty, metrics.best_stoch.penalty),
            "Key open task: improve stochastic branch while preserving nominal fidelity floor.",
            "This concludes the report.",
        ],
        size=18,
    )
    _add_picture(s, ASSET_DIR / "conclusion_summary.png", 7.85, 1.35, w=4.8)

    prs.save(str(PPT_PATH_DETAILED))
    shutil.copy2(PPT_PATH_DETAILED, PPT_PATH_ALIAS)
    return len(prs.slides)


def main() -> None:
    metrics = read_metrics()
    generate_assets(metrics)
    slide_count = build_presentation(metrics)

    print("Wrote detailed PPT:", PPT_PATH_DETAILED)
    print("Updated alias PPT:", PPT_PATH_ALIAS)
    print("Static best:", metrics.best_static.tag, "p=%.3f score=%.6f" % (metrics.best_static.penalty, metrics.best_static.score))
    print("Stochastic best:", metrics.best_stoch.tag, "p=%.3f score=%.6f" % (metrics.best_stoch.penalty, metrics.best_stoch.score))
    print("Slides:", slide_count)


if __name__ == "__main__":
    main()
